from django.shortcuts import render, redirect, get_object_or_404
from django.views.generic import ListView, CreateView, UpdateView, DeleteView, DetailView, TemplateView
from django.urls import reverse_lazy
from django.contrib import messages
from django.http import JsonResponse, HttpResponse, HttpResponseRedirect
from django.contrib.auth.mixins import LoginRequiredMixin, PermissionRequiredMixin
from django.db import transaction
from django.core.paginator import Paginator
from django.utils import timezone
from django.views.decorators.csrf import csrf_exempt
from django.utils.decorators import method_decorator
import json

from .services.workflow_service import WorkflowService

from .models import (
    AIModelConfig,
    AIWorkflow,
    WorkflowNode,
    WorkflowConnection,
    AIWorkflowExecution,
    AIChat,
    AIChatMessage,
    AIKnowledgeBase,
    AIKnowledgeItem,
    AIKnowledgeVector,
    AISalesStrategy,
    AIIntentRecognition,
    AIEmotionAnalysis,
    AIComplianceRule,
    AIActionTrigger,
    AILog
)
from .forms import (
    AIModelConfigForm,
    AIWorkflowForm,
    WorkflowNodeForm,
    WorkflowConnectionForm,
    AIKnowledgeBaseForm,
    AIKnowledgeItemForm,
    AISalesStrategyForm,
    AIIntentRecognitionForm,
    AIEmotionAnalysisForm,
    AIComplianceRuleForm,
    AIActionTriggerForm
)
from .utils.ai_config_manager import get_ai_config_manager
from .utils.ai_client import AIClient


# AI模型配置视图
class AIModelConfigListView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AIModelConfig
    template_name = 'ai/model_config_list.html'
    context_object_name = 'model_configs'
    permission_required = 'ai.view_aimodelconfig'
    paginate_by = 10
    
    def get_queryset(self):
        return AIModelConfig.objects.order_by('-created_at')
    
    def get(self, request, *args, **kwargs):
        # 检查是否为AJAX请求
        # 只通过X-Requested-With头判断，避免accepts导致的问题
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            # 处理AJAX请求，返回JSON数据
            queryset = self.get_queryset()
            paginator = Paginator(queryset, self.paginate_by)
            
            page = request.GET.get('page')
            objects = paginator.get_page(page)
            
            # 构造Layui表格需要的数据格式
            data = {
                "code": 0,
                "msg": "",
                "count": paginator.count,
                "data": [
                    {
                        "id": obj.id,
                        "name": obj.name,
                        "provider": obj.provider,
                        "model_type": obj.model_type,
                        "api_key": obj.api_key,
                        "api_base": obj.api_base,
                        "is_active": obj.is_active,
                        "created_at": obj.created_at.strftime('%Y-%m-%d %H:%M:%S')
                    } for obj in objects
                ]
            }
            
            return JsonResponse(data)
        else:
            # 处理HTML请求，返回完整页面
            return super().get(request, *args, **kwargs)


class AIModelConfigCreateView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    model = AIModelConfig
    form_class = AIModelConfigForm
    template_name = 'ai/model_config_form.html'
    permission_required = 'ai.add_aimodelconfig'
    success_url = reverse_lazy('ai:model_config_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, 'AI模型配置创建成功')
        return response


class AIModelConfigUpdateView(LoginRequiredMixin, PermissionRequiredMixin, UpdateView):
    model = AIModelConfig
    form_class = AIModelConfigForm
    template_name = 'ai/model_config_form.html'
    permission_required = 'ai.change_aimodelconfig'
    success_url = reverse_lazy('ai:model_config_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, 'AI模型配置更新成功')
        return response


class AIModelConfigDeleteView(LoginRequiredMixin, PermissionRequiredMixin, DeleteView):
    model = AIModelConfig
    template_name = 'ai/model_config_confirm_delete.html'
    permission_required = 'ai.delete_aimodelconfig'
    success_url = reverse_lazy('ai:model_config_list')
    
    def delete(self, request, *args, **kwargs):
        response = super().delete(request, *args, **kwargs)
        messages.success(self.request, 'AI模型配置删除成功')
        return response


class AIModelConfigDetailView(LoginRequiredMixin, PermissionRequiredMixin, DetailView):
    model = AIModelConfig
    template_name = 'ai/model_config_detail.html'
    context_object_name = 'model_config'
    permission_required = 'ai.view_aimodelconfig'


class AIModelConfigValidateView(LoginRequiredMixin, PermissionRequiredMixin, DetailView):
    model = AIModelConfig
    permission_required = 'ai.view_aimodelconfig'
    
    def get(self, request, *args, **kwargs):
        return self.validate_connection()
    
    def post(self, request, *args, **kwargs):
        return self.validate_connection()
    
    def validate_connection(self):
        model_config = self.get_object()
        try:
            # 记录详细的调试信息
            import logging
            logger = logging.getLogger(__name__)
            logger.info(f"测试AI模型连接 - 模型ID: {model_config.id}")
            logger.info(f"模型配置 - 提供商: {model_config.provider}, 基础URL: {model_config.api_base}, 模型名称: {model_config.model_name}")
            logger.info(f"API密钥: {'***' if model_config.api_key else '未配置'}")
            
            # 直接使用model_config_id参数实例化AIClient
            client = AIClient(model_config_id=model_config.id)
            # 测试连接，使用chat_completion方法发送简单消息
            test_message = [{"role": "user", "content": "你好，这是一个连接测试。"}]
            
            # 记录测试消息
            logger.info(f"测试消息: {test_message}")
            
            response = client.client.chat_completion(test_message)
            logger.info(f"AI模型连接成功 - 响应: {response[:50]}...")
            return JsonResponse({
                'status': 'success', 
                'message': '连接成功', 
                'result': response[:50] + '...',
                'details': {
                    'provider': model_config.provider,
                    'base_url': model_config.api_base,
                    'model_name': model_config.model_name,
                    'test_message': test_message
                }
            })
        except Exception as e:
            # 记录详细的错误信息
            import logging
            import traceback
            logger = logging.getLogger(__name__)
            logger.error(f"AI模型连接失败 - 模型ID: {model_config.id}")
            logger.error(f"错误类型: {type(e).__name__}")
            logger.error(f"错误详情: {str(e)}")
            logger.error(f"完整错误堆栈: {traceback.format_exc()}")
            logger.error(f"模型配置 - 提供商: {model_config.provider}, 基础URL: {model_config.api_base}, 模型名称: {model_config.model_name}")
            
            return JsonResponse({
                'status': 'error', 
                'message': f'连接失败: {str(e)}',
                'details': {
                    'provider': model_config.provider,
                    'base_url': model_config.api_base,
                    'model_name': model_config.model_name,
                    'error_type': type(e).__name__,
                    'error_message': str(e),
                    'suggestion': self._get_error_suggestion(type(e).__name__, model_config)
                }
            })
    
    def _get_error_suggestion(self, error_type, model_config):
        """根据错误类型提供修复建议"""
        suggestions = {
            'ConnectionError': '请检查网络连接是否正常，以及API地址是否正确',
            'TimeoutError': '请求超时，请检查网络连接或API地址是否正确',
            'HTTPError': f'HTTP请求失败，请检查API地址是否正确。当前配置的地址是: {model_config.api_base}',
            'AIClientError': 'AI客户端错误，请检查API密钥和API地址是否正确',
            'KeyError': 'API响应格式错误，请检查API地址是否正确',
            'ValueError': '参数错误，请检查模型配置是否正确',
            'ImportError': '缺少依赖库，请安装相关依赖',
            'AttributeError': '代码错误，请联系开发人员',
            'TypeError': '类型错误，请检查API地址格式是否正确'
        }
        
        return suggestions.get(error_type, '请检查模型配置是否正确，特别是API地址和API密钥')


class ListAvailableAIProvidersView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    permission_required = 'ai.view_aimodelconfig'
    
    def get(self, request, *args, **kwargs):
        providers = AIModelConfig.PROVIDERS
        return JsonResponse({'providers': [{'value': p[0], 'label': p[1]} for p in providers]})


# AI工作流视图
class AIWorkflowListView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AIWorkflow
    template_name = 'ai/workflow_list.html'
    context_object_name = 'workflows'
    permission_required = 'ai.view_aiworkflow'
    paginate_by = 10
    
    def get_queryset(self):
        return AIWorkflow.objects.order_by('-created_at')
    
    def get(self, request, *args, **kwargs):
        # 检查是否为AJAX请求
        # 只通过X-Requested-With头判断，避免accepts导致的问题
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            # 处理AJAX请求，返回JSON数据
            queryset = self.get_queryset()
            paginator = Paginator(queryset, self.paginate_by)
            
            page = request.GET.get('page')
            objects = paginator.get_page(page)
            
            # 构造Layui表格需要的数据格式
            data = {
                "code": 0,
                "msg": "",
                "count": paginator.count,
                "data": [
                    {
                        "id": obj.id,
                        "name": obj.name,
                        "status": obj.status,
                        "created_by__name": obj.created_by.name if obj.created_by and hasattr(obj.created_by, 'name') else '',
                        "created_at": obj.created_at.strftime('%Y-%m-%d %H:%M:%S')
                    } for obj in objects
                ]
            }
            
            return JsonResponse(data)
        else:
            # 处理HTML请求，返回完整页面
            return super().get(request, *args, **kwargs)


class AIWorkflowCreateView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    model = AIWorkflow
    form_class = AIWorkflowForm
    template_name = 'ai/workflow_form.html'
    permission_required = 'ai.add_aiworkflow'
    success_url = reverse_lazy('ai:workflow_list')
    
    def form_valid(self, form):
        form.instance.owner = self.request.user
        form.instance.created_by = self.request.user
        response = super().form_valid(form)
        messages.success(self.request, 'AI工作流创建成功')
        return response


class AIWorkflowUpdateView(LoginRequiredMixin, PermissionRequiredMixin, UpdateView):
    model = AIWorkflow
    form_class = AIWorkflowForm
    template_name = 'ai/workflow_form.html'
    permission_required = 'ai.change_aiworkflow'
    success_url = reverse_lazy('ai:workflow_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, 'AI工作流更新成功')
        return response


@method_decorator(csrf_exempt, name='dispatch')
class AIWorkflowDeleteView(LoginRequiredMixin, PermissionRequiredMixin, DeleteView):
    model = AIWorkflow
    permission_required = 'ai.delete_aiworkflow'
    
    def delete(self, request, *args, **kwargs):
        workflow = self.get_object()
        workflow_name = workflow.name
        
        try:
            workflow.delete()
            return JsonResponse({
                'success': True,
                'message': f'工作流 "{workflow_name}" 删除成功'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'message': f'删除失败：{str(e)}'
            }, status=500)
    
    def post(self, request, *args, **kwargs):
        return self.delete(request, *args, **kwargs)


class AIWorkflowDesignerView(LoginRequiredMixin, PermissionRequiredMixin, DetailView):
    model = AIWorkflow
    template_name = 'ai/workflow_designer.html'
    permission_required = 'ai.change_aiworkflow'
    context_object_name = 'workflow'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['nodes'] = self.object.nodes.all()
        context['connections'] = self.object.connections.all()
        return context


class AIWorkflowPublishView(LoginRequiredMixin, PermissionRequiredMixin, UpdateView):
    model = AIWorkflow
    fields = ['status']
    permission_required = 'ai.change_aiworkflow'
    
    def post(self, request, *args, **kwargs):
        workflow = self.get_object()
        
        if workflow.status == 'published':
            return JsonResponse({
                'status': 'error',
                'message': '工作流已经发布过了'
            })
        
        workflow.status = 'published'
        workflow.save()
        
        return JsonResponse({
            'status': 'success',
            'message': '发布成功'
        })


class AIWorkflowDetailView(LoginRequiredMixin, PermissionRequiredMixin, DetailView):
    model = AIWorkflow
    template_name = 'ai/workflow_detail.html'
    permission_required = 'ai.view_aiworkflow'
    context_object_name = 'workflow'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['nodes'] = self.object.nodes.all()
        context['connections'] = self.object.connections.all()
        return context


class AIWorkflowJsonDetailView(LoginRequiredMixin, DetailView):
    """工作流JSON详情API，用于工作流设计器数据加载"""
    model = AIWorkflow
    
    def get(self, request, *args, **kwargs):
        try:
            workflow = self.get_object()
            
            if not workflow.is_public:
                if not request.user.is_authenticated:
                    return JsonResponse({
                        'status': 'error',
                        'message': '请先登录'
                    }, status=401)
            
            nodes = workflow.nodes.all()
            connections = workflow.connections.all()
            
            def parse_config(config):
                if config is None:
                    return {}
                if isinstance(config, str):
                    try:
                        return json.loads(config)
                    except:
                        return {}
                return config if isinstance(config, dict) else {}
            
            nodes_data = []
            for node in nodes:
                node_data = {
                    'id': str(node.id),
                    'name': node.name,
                    'type': node.node_type or 'basic',
                    'x': int(node.position_x) if node.position_x else 100,
                    'y': int(node.position_y) if node.position_y else 200,
                    'config': parse_config(node.config)
                }
                nodes_data.append(node_data)
            
            connections_data = []
            for conn in connections:
                conn_data = {
                    'id': str(conn.id),
                    'source': str(conn.source_node_id) if conn.source_node_id else None,
                    'target': str(conn.target_node_id) if conn.target_node_id else None,
                    'source_handle': conn.source_handle or 'output',
                    'target_handle': conn.target_handle or 'input',
                    'config': parse_config(conn.config)
                }
                connections_data.append(conn_data)
            
            return JsonResponse({
                'status': 'success',
                'workflow': {
                    'id': str(workflow.id),
                    'name': workflow.name,
                    'description': workflow.description or '',
                    'status': workflow.status or 'draft',
                    'nodes': nodes_data,
                    'connections': connections_data
                }
            })
        except Exception as e:
            return JsonResponse({
                'status': 'error',
                'message': str(e)
            }, status=500)


class AIWorkflowExecuteView(LoginRequiredMixin, PermissionRequiredMixin, DetailView):
    model = AIWorkflow
    permission_required = 'ai.change_aiworkflow'
    
    def post(self, request, *args, **kwargs):
        import json
        
        workflow = self.get_object()
        
        try:
            # 尝试从JSON body获取输入数据
            if request.content_type == 'application/json':
                try:
                    body_data = json.loads(request.body)
                    input_data = body_data.get('input_data', {})
                except json.JSONDecodeError:
                    input_data = {}
            else:
                input_data = request.POST.dict()
            
            service = WorkflowService()
            execution = service.execute_workflow(
                workflow_id=str(workflow.id),
                user=request.user,
                input_data=input_data
            )
            return JsonResponse({
                'status': 'success', 
                'execution_id': str(execution.id),
                'message': '工作流执行已启动'
            })
        except Exception as e:
            return JsonResponse({
                'status': 'error', 
                'message': str(e)
            }, status=400)


class AIWorkflowParametersView(LoginRequiredMixin, PermissionRequiredMixin, DetailView):
    model = AIWorkflow
    permission_required = 'ai.change_aiworkflow'
    
    def get(self, request, *args, **kwargs):
        workflow = self.get_object()
        
        try:
            # 从WorkflowVariable获取定义的参数
            from apps.ai.models import WorkflowVariable
            
            variables = workflow.variables.all()
            parameters = []
            
            for var in variables:
                parameters.append({
                    'name': var.name,
                    'data_type': var.data_type,
                    'default_value': var.default_value,
                    'description': var.description,
                    'is_required': var.is_required
                })
            
            # 如果没有定义变量，返回空数组
            return JsonResponse({
                'status': 'success',
                'parameters': parameters,
                'workflow_name': workflow.name
            })
            
        except Exception as e:
            return JsonResponse({
                'status': 'error',
                'message': str(e)
            }, status=400)


# AI聊天视图
class AIChatView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AIChat
    template_name = 'ai/chat.html'
    permission_required = 'ai.view_aichat'
    
    def get_queryset(self):
        return AIChat.objects.filter(user=self.request.user).order_by('-updated_at')
    
    def get(self, request, *args, **kwargs):
        # 处理JSON请求，返回聊天会话列表
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            queryset = self.get_queryset()
            chats = [
                {
                    'id': chat.id,
                    'title': chat.title,
                    'created_at': chat.created_at.strftime('%Y-%m-%d %H:%M:%S'),
                    'updated_at': chat.updated_at.strftime('%Y-%m-%d %H:%M:%S')
                } for chat in queryset
            ]
            return JsonResponse(chats, safe=False)
        return super().get(request, *args, **kwargs)


class AIChatDetailView(LoginRequiredMixin, PermissionRequiredMixin, DetailView):
    model = AIChat
    template_name = 'ai/chat_detail.html'
    permission_required = 'ai.view_aichat'
    context_object_name = 'chat'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['messages'] = self.object.messages.order_by('created_at')
        return context
    
    def get(self, request, *args, **kwargs):
        # 处理JSON请求，返回聊天会话数据
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            chat = self.get_object()
            messages = chat.messages.order_by('created_at')
            
            return JsonResponse({
                'id': chat.id,
                'title': chat.title,
                'messages': [
                    {
                        'id': msg.id,
                        'role': msg.role,
                        'content': msg.content,
                        'created_at': msg.created_at.strftime('%Y-%m-%d %H:%M:%S')
                    } for msg in messages
                ]
            })
        return super().get(request, *args, **kwargs)


class AIChatDeleteView(LoginRequiredMixin, PermissionRequiredMixin, DeleteView):
    model = AIChat
    permission_required = 'ai.delete_aichat'
    success_url = reverse_lazy('ai:chat')
    
    def delete(self, request, *args, **kwargs):
        chat = self.get_object()
        chat.delete()
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            return JsonResponse({'status': 'success', 'message': '聊天会话已删除'})
        return super().delete(request, *args, **kwargs)


class AIChatCreateView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    model = AIChat
    fields = ['title']
    permission_required = 'ai.add_aichat'
    
    def form_valid(self, form):
        form.instance.user = self.request.user
        form.instance.session_id = f"chat_{self.request.user.id}_{timezone.now().timestamp()}"
        chat = form.save()
        return JsonResponse({'status': 'success', 'chat_id': chat.id, 'chat_title': chat.title})


class AIChatMessageCreateView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    model = AIChatMessage
    fields = ['content']
    permission_required = 'ai.add_aichatmessage'
    
    def form_valid(self, form):
        chat_id = self.request.POST.get('chat_id')
        chat = get_object_or_404(AIChat, id=chat_id, user=self.request.user)
        message_content = form.cleaned_data['content']
        
        # 创建用户消息
        user_message = AIChatMessage.objects.create(
            chat=chat,
            role='user',
            content=message_content
        )
        
        # 生成AI回复
        try:
            # 1. 调用意图识别服务
            from apps.ai.services.intent_recognition_service import intent_recognition_service
            intent_result = intent_recognition_service.recognize_intent(self.request.user, message_content)
            
            # 2. 获取意图特定处理器
            intent = intent_result['intent']
            
            # 直接处理意图，避免handler绑定问题
            if intent == intent_recognition_service.INTENT_TYPES['DATA_QUERY']:
                # 3. 执行数据查询服务
                from apps.ai.services.query_service import query_service
                result = query_service.process_query(self.request.user, message_content)
                
                if result['success']:
                    ai_response = result['result']
                else:
                    ai_response = result.get('message', '抱歉，我无法处理您的请求')
            elif intent == intent_recognition_service.INTENT_TYPES['KNOWLEDGE_BASE']:
                from apps.ai.services.rag_service import rag_service
                ai_response = rag_service.generate_response(self.request.user, message_content)
            elif intent == intent_recognition_service.INTENT_TYPES['AI_CHAT']:
                # 5. 处理纯AI对话意图
                ai_response = intent_recognition_service._handle_ai_chat(self.request.user, message_content)['result']
            else:
                # 6. 如果没有特定处理器，使用通用AI聊天
                ai_response = intent_result.get('fallback_options', [{}])[0].get('text', '抱歉，我无法理解您的请求')
            
            # 5. 创建AI回复消息
            ai_message = AIChatMessage.objects.create(
                chat=chat,
                role='assistant',
                content=ai_response
            )
            
            return JsonResponse({
                'status': 'success',
                'user_message': user_message.content,
                'ai_message': ai_message.content,
                'intent': intent,
                'confidence': intent_result['confidence']
            })
        except Exception as e:
            logger.error(f'AI生成失败: {str(e)}')
            return JsonResponse({'status': 'error', 'message': f'AI生成失败: {str(e)}'})


# AI知识库视图
class AIKnowledgeBaseListView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AIKnowledgeBase
    template_name = 'ai/knowledge_base_list.html'
    context_object_name = 'knowledge_bases'
    permission_required = 'ai.view_aiknowledgebase'
    paginate_by = 10
    
    def get_queryset(self):
        return AIKnowledgeBase.objects.order_by('-created_at')
    
    def get(self, request, *args, **kwargs):
        # 检查是否为AJAX请求
        # 只通过X-Requested-With头判断，避免accepts导致的问题
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            # 处理AJAX请求，返回JSON数据
            queryset = self.get_queryset()
            paginator = Paginator(queryset, self.paginate_by)
            
            page = request.GET.get('page')
            objects = paginator.get_page(page)
            
            # 构造Layui表格需要的数据格式
            data = {
                "code": 0,
                "msg": "",
                "count": paginator.count,
                "data": [
                    {
                        "id": obj.id,
                        "name": obj.name,
                        "description": obj.description,
                        "status": obj.status,
                        "creator__name": obj.creator.name if obj.creator and hasattr(obj.creator, 'name') else '',
                        "created_at": obj.created_at.strftime('%Y-%m-%d %H:%M:%S')
                    } for obj in objects
                ]
            }
            
            return JsonResponse(data)
        else:
            # 处理HTML请求，返回完整页面
            return super().get(request, *args, **kwargs)


class AIKnowledgeBaseCreateView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    model = AIKnowledgeBase
    form_class = AIKnowledgeBaseForm
    template_name = 'ai/knowledge_base_form.html'
    permission_required = 'ai.add_aiknowledgebase'
    success_url = reverse_lazy('ai:knowledge_base_list')
    
    def form_valid(self, form):
        form.instance.creator = self.request.user
        response = super().form_valid(form)
        messages.success(self.request, 'AI知识库创建成功')
        return response


class AIKnowledgeBaseUpdateView(LoginRequiredMixin, PermissionRequiredMixin, UpdateView):
    model = AIKnowledgeBase
    form_class = AIKnowledgeBaseForm
    template_name = 'ai/knowledge_base_form.html'
    permission_required = 'ai.change_aiknowledgebase'
    success_url = reverse_lazy('ai:knowledge_base_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, 'AI知识库更新成功')
        return response


class AIKnowledgeBaseDeleteView(LoginRequiredMixin, PermissionRequiredMixin, DeleteView):
    model = AIKnowledgeBase
    template_name = 'ai/knowledge_base_confirm_delete.html'
    permission_required = 'ai.delete_aiknowledgebase'
    success_url = reverse_lazy('ai:knowledge_base_list')
    
    def delete(self, request, *args, **kwargs):
        response = super().delete(request, *args, **kwargs)
        messages.success(self.request, 'AI知识库删除成功')
        return response


class AIKnowledgeBaseDetailView(LoginRequiredMixin, PermissionRequiredMixin, DetailView):
    model = AIKnowledgeBase
    template_name = 'ai/knowledge_base_detail.html'
    permission_required = 'ai.view_aiknowledgebase'
    context_object_name = 'knowledge_base'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['knowledge_items'] = self.object.items.all().order_by('-created_at')
        return context


# AI知识条目视图
class AIKnowledgeItemListView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AIKnowledgeItem
    template_name = 'ai/knowledge_item_list.html'
    context_object_name = 'knowledge_items'
    permission_required = 'ai.view_aiknowledgeitem'
    paginate_by = 10
    
    def get_queryset(self):
        knowledge_base_id = self.request.GET.get('knowledge_base')
        if knowledge_base_id:
            return AIKnowledgeItem.objects.filter(knowledge_base_id=knowledge_base_id).order_by('-created_at')
        return AIKnowledgeItem.objects.order_by('-created_at')
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['knowledge_bases'] = AIKnowledgeBase.objects.filter(status='published')
        return context
    
    def get(self, request, *args, **kwargs):
        # 检查是否为AJAX请求
        # 只通过X-Requested-With头判断，避免accepts导致的问题
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            # 处理AJAX请求，返回JSON数据
            queryset = self.get_queryset()
            paginator = Paginator(queryset, self.paginate_by)
            
            page = request.GET.get('page')
            objects = paginator.get_page(page)
            
            # 构造Layui表格需要的数据格式
            data = {
                "code": 0,
                "msg": "",
                "count": paginator.count,
                "data": [
                    {
                        "id": obj.id,
                        "title": obj.title,
                        "knowledge_type": obj.knowledge_type,
                        "status": obj.status,
                        "knowledge_base__name": obj.knowledge_base.name if obj.knowledge_base and hasattr(obj.knowledge_base, 'name') else '',
                        "creator__name": obj.creator.name if obj.creator and hasattr(obj.creator, 'name') else '',
                        "created_at": obj.created_at.strftime('%Y-%m-%d %H:%M:%S'),
                        "has_file": True if obj.file else False,
                        "file_name": obj.file.name.split('/')[-1] if obj.file else '',
                        "file_type": obj.file_type if obj.file_type else ''
                    } for obj in objects
                ]
            }
            
            return JsonResponse(data)
        else:
            # 处理HTML请求，返回完整页面
            return super().get(request, *args, **kwargs)


class AIKnowledgeItemCreateView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    model = AIKnowledgeItem
    form_class = AIKnowledgeItemForm
    template_name = 'ai/knowledge_item_form.html'
    permission_required = 'ai.add_aiknowledgeitem'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        # 添加所有知识库到上下文
        context['knowledge_bases'] = AIKnowledgeBase.objects.all()
        return context
    
    def form_valid(self, form):
        try:
            # 获取表单数据
            form.instance.creator = self.request.user
            
            # 保存知识库条目
            knowledge_item = form.save()
            
            # 生成向量
            from apps.ai.services.vector_generation_service import vector_generation_service
            vector_generation_success = vector_generation_service.generate_vector_for_knowledge_item(knowledge_item.id)
            
            # 检查向量是否生成成功
            if not vector_generation_success:
                # 向量化失败，删除已保存的知识库条目
                knowledge_item.delete()
                # 返回错误信息
                if self.request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                    return JsonResponse({'status': 'error', 'message': '向量化失败，请检查内容并重试'})
                else:
                    messages.error(self.request, '向量化失败，请检查内容并重试')
                    return redirect('ai:knowledge_item_list')
            
            # 检查是否是AJAX请求
            if self.request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return JsonResponse({'status': 'success', 'message': '知识条目创建成功'})
            else:
                messages.success(self.request, '知识条目创建成功')
                return redirect('ai:knowledge_item_list')
        except Exception as e:
            # 记录错误日志
            import logging
            logger = logging.getLogger(__name__)
            logger.error(f"创建知识库条目失败: {str(e)}")
            
            # 返回错误信息
            error_msg = f"操作失败: {str(e)}"
            if self.request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return JsonResponse({'status': 'error', 'message': error_msg})
            else:
                messages.error(self.request, error_msg)
                return self.form_invalid(form)


class AIKnowledgeItemUpdateView(LoginRequiredMixin, PermissionRequiredMixin, UpdateView):
    model = AIKnowledgeItem
    form_class = AIKnowledgeItemForm
    template_name = 'ai/knowledge_item_form.html'
    permission_required = 'ai.change_aiknowledgeitem'
    success_url = reverse_lazy('ai:knowledge_item_list')
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        # 添加所有知识库到上下文
        context['knowledge_bases'] = AIKnowledgeBase.objects.all()
        return context
    
    def form_valid(self, form):
        try:
            # 保存知识库条目更新
            knowledge_item = form.save()
            
            # 生成向量
            from apps.ai.services.vector_generation_service import vector_generation_service
            vector_generation_success = vector_generation_service.generate_vector_for_knowledge_item(knowledge_item.id)
            
            # 检查向量是否生成成功
            if not vector_generation_success:
                # 向量化失败，返回错误信息
                if self.request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                    return JsonResponse({'status': 'error', 'message': '向量化失败，知识条目已更新但向量未更新'})
                else:
                    messages.error(self.request, '向量化失败，知识条目已更新但向量未更新')
                    return redirect('ai:knowledge_item_list')
            
            # 检查是否是AJAX请求
            if self.request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return JsonResponse({'status': 'success', 'message': '知识条目更新成功'})
            else:
                messages.success(self.request, '知识条目更新成功')
                return redirect('ai:knowledge_item_list')
        except Exception as e:
            # 记录错误日志
            import logging
            logger = logging.getLogger(__name__)
            logger.error(f"更新知识库条目失败: {str(e)}")
            
            # 返回错误信息
            error_msg = f"操作失败: {str(e)}"
            if self.request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return JsonResponse({'status': 'error', 'message': error_msg})
            else:
                messages.error(self.request, error_msg)
                return self.form_invalid(form)


class AIKnowledgeItemDeleteView(LoginRequiredMixin, PermissionRequiredMixin, DeleteView):
    model = AIKnowledgeItem
    template_name = 'ai/knowledge_item_confirm_delete.html'
    permission_required = 'ai.delete_aiknowledgeitem'
    success_url = reverse_lazy('ai:knowledge_item_list')
    
    def delete(self, request, *args, **kwargs):
        # 执行删除操作
        super().delete(request, *args, **kwargs)
        messages.success(self.request, '知识条目删除成功')
        
        # 检查是否是AJAX请求
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            # 返回JSON响应
            return JsonResponse({'status': 'success', 'message': '知识条目删除成功'})
        else:
            # 返回重定向响应
            return HttpResponseRedirect(self.success_url)
    
    def post(self, request, *args, **kwargs):
        # 处理POST请求，调用delete方法
        return self.delete(request, *args, **kwargs)


class AIKnowledgeItemDetailView(LoginRequiredMixin, PermissionRequiredMixin, DetailView):
    model = AIKnowledgeItem
    template_name = 'ai/knowledge_item_detail.html'
    permission_required = 'ai.view_aiknowledgeitem'
    context_object_name = 'knowledge_item'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        
        # 获取知识条目的向量信息
        knowledge_item = self.get_object()
        try:
            vector_record = AIKnowledgeVector.objects.get(knowledge_item=knowledge_item)
            context['vector_info'] = {
                'exists': True,
                'dimension': vector_record.dimension,
                'created_at': vector_record.created_at,
                'updated_at': vector_record.updated_at,
                'vector_size': len(vector_record.vector) if vector_record.vector else 0
            }
        except AIKnowledgeVector.DoesNotExist:
            context['vector_info'] = {
                'exists': False
            }
        
        return context


class AIKnowledgeSearchView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AIKnowledgeItem
    template_name = 'ai/knowledge_search.html'
    permission_required = 'ai.view_aiknowledgeitem'
    
    def get_queryset(self):
        query = self.request.GET.get('q', '')
        if not query:
            return AIKnowledgeItem.objects.none()
        
        return AIKnowledgeItem.objects.filter(
            content__icontains=query,
            status='published'
        ).order_by('-created_at')
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['query'] = self.request.GET.get('q', '')
        return context


# AI销售策略视图
class AISalesStrategyListView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AISalesStrategy
    template_name = 'ai/sales_strategy_list.html'
    context_object_name = 'sales_strategies'
    permission_required = 'ai.view_aisalesstrategy'
    paginate_by = 10
    
    def get_queryset(self):
        return AISalesStrategy.objects.order_by('-created_at')


class AISalesStrategyCreateView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    model = AISalesStrategy
    form_class = AISalesStrategyForm
    template_name = 'ai/sales_strategy_form.html'
    permission_required = 'ai.add_aisalesstrategy'
    success_url = reverse_lazy('ai:sales_strategy_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, '销售策略创建成功')
        return response


class AISalesStrategyUpdateView(LoginRequiredMixin, PermissionRequiredMixin, UpdateView):
    model = AISalesStrategy
    form_class = AISalesStrategyForm
    template_name = 'ai/sales_strategy_form.html'
    permission_required = 'ai.change_aisalesstrategy'
    success_url = reverse_lazy('ai:sales_strategy_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, '销售策略更新成功')
        return response


class AISalesStrategyDeleteView(LoginRequiredMixin, PermissionRequiredMixin, DeleteView):
    model = AISalesStrategy
    template_name = 'ai/sales_strategy_confirm_delete.html'
    permission_required = 'ai.delete_aisalesstrategy'
    success_url = reverse_lazy('ai:sales_strategy_list')
    
    def delete(self, request, *args, **kwargs):
        response = super().delete(request, *args, **kwargs)
        messages.success(self.request, '销售策略删除成功')
        return response


# AI意图识别视图
class AIIntentRecognitionListView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AIIntentRecognition
    template_name = 'ai/intent_recognition_list.html'
    context_object_name = 'intent_recognitions'
    permission_required = 'ai.view_aiintentrecognition'
    paginate_by = 10
    
    def get_queryset(self):
        return AIIntentRecognition.objects.order_by('-created_at')


class AIIntentRecognitionCreateView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    model = AIIntentRecognition
    form_class = AIIntentRecognitionForm
    template_name = 'ai/intent_recognition_form.html'
    permission_required = 'ai.add_aiintentrecognition'
    success_url = reverse_lazy('ai:intent_recognition_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, '意图识别规则创建成功')
        return response


class AIIntentRecognitionUpdateView(LoginRequiredMixin, PermissionRequiredMixin, UpdateView):
    model = AIIntentRecognition
    form_class = AIIntentRecognitionForm
    template_name = 'ai/intent_recognition_form.html'
    permission_required = 'ai.change_aiintentrecognition'
    success_url = reverse_lazy('ai:intent_recognition_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, '意图识别规则更新成功')
        return response


class AIIntentRecognitionDeleteView(LoginRequiredMixin, PermissionRequiredMixin, DeleteView):
    model = AIIntentRecognition
    template_name = 'ai/intent_recognition_confirm_delete.html'
    permission_required = 'ai.delete_aiintentrecognition'
    success_url = reverse_lazy('ai:intent_recognition_list')
    
    def delete(self, request, *args, **kwargs):
        response = super().delete(request, *args, **kwargs)
        messages.success(self.request, '意图识别规则删除成功')
        return response


# AI情绪分析视图
class AIEmotionAnalysisListView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AIEmotionAnalysis
    template_name = 'ai/emotion_analysis_list.html'
    context_object_name = 'emotion_analyses'
    permission_required = 'ai.view_aiemotionanalysis'
    paginate_by = 10
    
    def get_queryset(self):
        return AIEmotionAnalysis.objects.order_by('-created_at')


class AIEmotionAnalysisCreateView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    model = AIEmotionAnalysis
    form_class = AIEmotionAnalysisForm
    template_name = 'ai/emotion_analysis_form.html'
    permission_required = 'ai.add_aiemotionanalysis'
    success_url = reverse_lazy('ai:emotion_analysis_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, '情绪分析规则创建成功')
        return response


class AIEmotionAnalysisUpdateView(LoginRequiredMixin, PermissionRequiredMixin, UpdateView):
    model = AIEmotionAnalysis
    form_class = AIEmotionAnalysisForm
    template_name = 'ai/emotion_analysis_form.html'
    permission_required = 'ai.change_aiemotionanalysis'
    success_url = reverse_lazy('ai:emotion_analysis_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, '情绪分析规则更新成功')
        return response


class AIEmotionAnalysisDeleteView(LoginRequiredMixin, PermissionRequiredMixin, DeleteView):
    model = AIEmotionAnalysis
    template_name = 'ai/emotion_analysis_confirm_delete.html'
    permission_required = 'ai.delete_aiemotionanalysis'
    success_url = reverse_lazy('ai:emotion_analysis_list')
    
    def delete(self, request, *args, **kwargs):
        response = super().delete(request, *args, **kwargs)
        messages.success(self.request, '情绪分析规则删除成功')
        return response


# AI合规规则视图
class AIComplianceRuleListView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AIComplianceRule
    template_name = 'ai/compliance_rule_list.html'
    context_object_name = 'compliance_rules'
    permission_required = 'ai.view_aicompliancerule'
    paginate_by = 10
    
    def get_queryset(self):
        return AIComplianceRule.objects.order_by('-created_at')


class AIComplianceRuleCreateView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    model = AIComplianceRule
    form_class = AIComplianceRuleForm
    template_name = 'ai/compliance_rule_form.html'
    permission_required = 'ai.add_aicompliancerule'
    success_url = reverse_lazy('ai:compliance_rule_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, '合规规则创建成功')
        return response


class AIComplianceRuleUpdateView(LoginRequiredMixin, PermissionRequiredMixin, UpdateView):
    model = AIComplianceRule
    form_class = AIComplianceRuleForm
    template_name = 'ai/compliance_rule_form.html'
    permission_required = 'ai.change_aicompliancerule'
    success_url = reverse_lazy('ai:compliance_rule_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, '合规规则更新成功')
        return response


class AIComplianceRuleDeleteView(LoginRequiredMixin, PermissionRequiredMixin, DeleteView):
    model = AIComplianceRule
    template_name = 'ai/compliance_rule_confirm_delete.html'
    permission_required = 'ai.delete_aicompliancerule'
    success_url = reverse_lazy('ai:compliance_rule_list')
    
    def delete(self, request, *args, **kwargs):
        response = super().delete(request, *args, **kwargs)
        messages.success(self.request, '合规规则删除成功')
        return response


# AI自动行动触发视图
class AIActionTriggerListView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AIActionTrigger
    template_name = 'ai/action_trigger_list.html'
    context_object_name = 'action_triggers'
    permission_required = 'ai.view_aiactiontrigger'
    paginate_by = 10
    
    def get_queryset(self):
        return AIActionTrigger.objects.order_by('-created_at')


class AIActionTriggerCreateView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    model = AIActionTrigger
    form_class = AIActionTriggerForm
    template_name = 'ai/action_trigger_form.html'
    permission_required = 'ai.add_aiactiontrigger'
    success_url = reverse_lazy('ai:action_trigger_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, '自动行动触发创建成功')
        return response


class AIActionTriggerUpdateView(LoginRequiredMixin, PermissionRequiredMixin, UpdateView):
    model = AIActionTrigger
    form_class = AIActionTriggerForm
    template_name = 'ai/action_trigger_form.html'
    permission_required = 'ai.change_aiactiontrigger'
    success_url = reverse_lazy('ai:action_trigger_list')
    
    def form_valid(self, form):
        response = super().form_valid(form)
        messages.success(self.request, '自动行动触发更新成功')
        return response


class AIActionTriggerDeleteView(LoginRequiredMixin, PermissionRequiredMixin, DeleteView):
    model = AIActionTrigger
    template_name = 'ai/action_trigger_confirm_delete.html'
    permission_required = 'ai.delete_aiactiontrigger'
    success_url = reverse_lazy('ai:action_trigger_list')
    
    def delete(self, request, *args, **kwargs):
        response = super().delete(request, *args, **kwargs)
        messages.success(self.request, '自动行动触发删除成功')
        return response


# AI仪表盘视图
class AIDashboardView(LoginRequiredMixin, PermissionRequiredMixin, TemplateView):
    template_name = 'ai/dashboard.html'
    permission_required = 'ai.view_aidashboard'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        # 添加仪表盘数据
        context['total_knowledge_bases'] = AIKnowledgeBase.objects.count()
        context['total_knowledge_items'] = AIKnowledgeItem.objects.count()
        context['total_sales_strategies'] = AISalesStrategy.objects.count()
        context['total_workflows'] = AIWorkflow.objects.count()
        context['total_model_configs'] = AIModelConfig.objects.count()
        return context


# AI操作日志视图
class AILogListView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    model = AILog
    template_name = 'ai/log_list.html'
    context_object_name = 'logs'
    permission_required = 'ai.view_ailog'
    paginate_by = 20
    
    def get_queryset(self):
        return AILog.objects.order_by('-created_at')
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['log_types'] = AILog.LOG_TYPES
        return context


# 文件解析视图
class ParseFileContentView(LoginRequiredMixin, PermissionRequiredMixin, CreateView):
    """文件内容解析视图"""
    permission_required = 'ai.add_aiknowledgeitem'
    
    def post(self, request, *args, **kwargs):
        """处理文件上传和解析"""
        try:
            if 'file' not in request.FILES:
                return JsonResponse({'status': 'error', 'message': '未找到上传的文件'})
            
            uploaded_file = request.FILES['file']
            file_name = uploaded_file.name
            file_extension = file_name.split('.')[-1].lower()
            
            # 读取文件内容
            file_content = uploaded_file.read()
            
            # 根据文件类型选择解析方法
            parsed_content = ''
            
            if file_extension == 'txt':
                # 解析TXT文件
                parsed_content = file_content.decode('utf-8', errors='ignore')
            elif file_extension in ['doc', 'docx']:
                # 解析Word文件
                try:
                    from docx import Document
                    from io import BytesIO
                    
                    doc = Document(BytesIO(file_content))
                    for para in doc.paragraphs:
                        parsed_content += para.text + '\n'
                except Exception as e:
                    return JsonResponse({'status': 'error', 'message': f'Word文件解析失败: {str(e)}'})
            elif file_extension in ['pdf']:
                # 解析PDF文件
                try:
                    import PyPDF2
                    from io import BytesIO
                    
                    pdf_reader = PyPDF2.PdfReader(BytesIO(file_content))
                    for page in pdf_reader.pages:
                        parsed_content += page.extract_text() + '\n'
                except Exception as e:
                    return JsonResponse({'status': 'error', 'message': f'PDF文件解析失败: {str(e)}'})
            elif file_extension in ['xls', 'xlsx']:
                # 解析Excel文件
                try:
                    import openpyxl
                    from io import BytesIO
                    
                    workbook = openpyxl.load_workbook(BytesIO(file_content))
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        parsed_content += f'=== {sheet_name} ===\n'
                        for row in sheet.iter_rows(values_only=True):
                            # 过滤掉空行
                            if any(cell is not None and cell != '' for cell in row):
                                row_content = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                                parsed_content += row_content + '\n'
                except Exception as e:
                    return JsonResponse({'status': 'error', 'message': f'Excel文件解析失败: {str(e)}'})
            elif file_extension in ['ppt', 'pptx']:
                # 解析PPT文件
                try:
                    from pptx import Presentation
                    from io import BytesIO
                    
                    presentation = Presentation(BytesIO(file_content))
                    for i, slide in enumerate(presentation.slides, 1):
                        parsed_content += f'=== 幻灯片 {i} ===\n'
                        for shape in slide.shapes:
                            if hasattr(shape, 'text'):
                                parsed_content += shape.text + '\n'
                except Exception as e:
                    return JsonResponse({'status': 'error', 'message': f'PPT文件解析失败: {str(e)}'})
            else:
                return JsonResponse({'status': 'error', 'message': f'不支持的文件类型: {file_extension}'})
            
            # 返回解析结果
            return JsonResponse({
                'status': 'success',
                'content': parsed_content,
                'file_name': file_name,
                'file_extension': file_extension
            })
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': f'文件解析失败: {str(e)}'})


# AI流式聊天视图
import logging

logger = logging.getLogger(__name__)

class AIChatStreamView(LoginRequiredMixin, CreateView):
    """AI流式聊天视图"""
    # 移除PermissionRequiredMixin，允许所有登录用户访问AI聊天功能
    # 权限检查通过RBAC中间件进行，而不是通过PermissionRequiredMixin
    
    def post(self, request, *args, **kwargs):
        """处理流式聊天请求"""
        import json
        
        try:
            # 获取请求数据
            if request.POST:
                # 表单数据
                data = request.POST
            else:
                # JSON数据
                try:
                    data = json.loads(request.body)
                except json.JSONDecodeError:
                    return JsonResponse({'status': 'error', 'message': '无效的JSON格式'})
            
            message = data.get('message', '')
            
            if not message:
                return JsonResponse({'status': 'error', 'message': '消息不能为空'})
            
            # 1. 调用意图识别服务
            from apps.ai.services.intent_recognition_service import intent_recognition_service
            intent_result = intent_recognition_service.recognize_intent(request.user, message)
            
            # 2. 获取意图特定处理器
            intent = intent_result['intent']
            
            # 直接处理意图，避免handler绑定问题
            if intent == intent_recognition_service.INTENT_TYPES['DATA_QUERY']:
                # 3. 执行数据查询服务
                from apps.ai.services.query_service import query_service
                result = query_service.process_query(request.user, message)
                
                if result['success']:
                    ai_response = result['result']
                else:
                    ai_response = result.get('message', '抱歉，我无法处理您的请求')
            elif intent == intent_recognition_service.INTENT_TYPES['KNOWLEDGE_BASE']:
                # 4. 处理知识库调用意图
                from apps.ai.services.rag_service import rag_service
                ai_response = rag_service.generate_response(request.user, message)
            elif intent == intent_recognition_service.INTENT_TYPES['AI_CHAT']:
                # 5. 处理纯AI对话意图
                chat_result = intent_recognition_service._handle_ai_chat(request.user, message)
                if chat_result.get('success'):
                    ai_response = chat_result.get('result', '')
                else:
                    ai_response = chat_result.get('message', '抱歉，AI服务暂时不可用')
            else:
                # 6. 处理意图识别失败情况
                fallback_options = intent_result.get('fallback_options', [])
                if fallback_options:
                    ai_response = "我不太确定您的需求，请选择以下选项：\n"
                    for i, option in enumerate(fallback_options, 1):
                        ai_response += f"{i}. {option['text']}\n"
                else:
                    ai_response = "抱歉，我无法理解您的请求，请尝试更明确的表述"
            
            try:
                # 5. 保存聊天记录
                chat = None
                chat_id = data.get('chat_id')
                
                if chat_id:
                    # 使用指定的聊天会话
                    try:
                        chat = AIChat.objects.get(id=chat_id, user=request.user)
                    except AIChat.DoesNotExist:
                        # 如果指定的聊天会话不存在，创建新的
                        chat = None
                
                if not chat:
                    # 创建新的聊天会话
                    chat, created = AIChat.objects.get_or_create(
                        user=request.user,
                        defaults={'title': f'聊天 {timezone.now().strftime("%Y-%m-%d %H:%M:%S")}'}
                    )
                
                # 保存用户消息
                AIChatMessage.objects.create(
                    chat=chat,
                    role='user',
                    content=message
                )
                
                # 保存AI回复
                AIChatMessage.objects.create(
                    chat=chat,
                    role='assistant',
                    content=ai_response
                )
            except Exception as e:
                logger.error(f'保存聊天记录失败: {str(e)}')
            
            # 6. 返回流式响应
            response = HttpResponse(self.generate_streaming_response(ai_response), content_type='text/event-stream')
            response['Cache-Control'] = 'no-cache'
            return response
            
        except Exception as e:
            logger.error(f'流式聊天请求失败: {str(e)}')
            error_message = f'抱歉，我暂时无法回答您的问题，请稍后再试。\n错误详情: {str(e)}'
            return HttpResponse(self.generate_streaming_response(error_message), content_type='text/event-stream')
    
    def generate_streaming_response(self, response_text):
        """生成流式响应"""
        # 模拟流式输出，逐字符发送
        for char in response_text:
            yield char
            import time
            time.sleep(0.01)  # 添加小延迟，模拟真实的流式输出


# AI工作流执行记录视图
class AIWorkflowExecutionListView(LoginRequiredMixin, PermissionRequiredMixin, ListView):
    """AI工作流执行记录列表视图"""
    model = AIWorkflowExecution
    template_name = 'ai/workflow_execution_list.html'
    permission_required = 'ai.view_aiworkflowexecution'
    context_object_name = 'executions'
    paginate_by = 10
    
    def get_queryset(self):
        """获取工作流执行记录，按开始时间倒序排列"""
        return AIWorkflowExecution.objects.order_by('-started_at')
    
    def get_context_data(self, **kwargs):
        """添加额外的上下文数据"""
        context = super().get_context_data(**kwargs)
        # 添加统计信息
        context['total_executions'] = AIWorkflowExecution.objects.count()
        context['running_executions'] = AIWorkflowExecution.objects.filter(status='running').count()
        context['completed_executions'] = AIWorkflowExecution.objects.filter(status='completed').count()
        context['failed_executions'] = AIWorkflowExecution.objects.filter(status='failed').count()
        return context


class AIWorkflowExecutionDetailView(LoginRequiredMixin, PermissionRequiredMixin, DetailView):
    """AI工作流执行记录详情视图"""
    model = AIWorkflowExecution
    template_name = 'ai/workflow_execution_detail.html'
    permission_required = 'ai.view_aiworkflowexecution'
    context_object_name = 'execution'
