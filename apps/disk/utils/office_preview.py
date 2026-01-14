import os
import re
import logging
import tempfile
import shutil
import signal
import threading
import unicodedata
from typing import Dict, Any, List, Optional
from pathlib import Path
from functools import wraps

# 设置日志 - 模块级别配置
logger = logging.getLogger('office_preview')

# 配置日志处理器（只执行一次）
def _setup_logger():
    """配置日志记录器"""
    if not logger.handlers:
        handler = logging.StreamHandler()
        formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
        handler.setFormatter(formatter)
        logger.addHandler(handler)
        logger.setLevel(logging.INFO)

_setup_logger()


class TextCleaner:
    """文本清理工具类 - 消除乱码和特殊字符"""
    
    CONTROL_CHARS = re.compile(
        '[\x00-\x08\x0b\x0c\x0e-\x1f\u200b-\u200f\u2028-\u202f\u2060-\u206f\ufeff]'
    )
    EXTRA_WHITESPACE = re.compile(r'[ \t]{2,}')
    MULTI_NEWLINES = re.compile(r'\n{3,}')
    
    @classmethod
    def clean_text(cls, text: str) -> str:
        if not text:
            return ''
        
        text = cls.CONTROL_CHARS.sub('', text)
        text = unicodedata.normalize('NFKC', text)
        text = cls.EXTRA_WHITESPACE.sub(' ', text)
        text = cls.MULTI_NEWLINES.sub('\n\n', text)
        text = '\n'.join(line.strip() for line in text.split('\n'))
        
        return text.strip()
    
    @classmethod
    def clean_run_text(cls, text: str) -> str:
        if not text:
            return ''
        
        text = cls.CONTROL_CHARS.sub('', text)
        text = unicodedata.normalize('NFKC', text)
        text = cls.EXTRA_WHITESPACE.sub(' ', text)
        
        return text.strip()


class PPTImageExtractor:
    """PPT图片提取器 - 提取并转换为Base64"""
    
    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.image_map = {}
        self.max_image_size = (1920, 1080)
        self.image_quality = 85
    
    def extract_all_images(self) -> dict:
        """提取所有图片"""
        import zipfile
        import base64
        from PIL import Image
        import io
        
        self.image_map = {}
        
        try:
            with zipfile.ZipFile(self.pptx_path, 'r') as zf:
                for name in zf.namelist():
                    if name.startswith('ppt/media/'):
                        self._extract_image(zf, name)
        except ImportError:
            logger.warning('Pillow未安装，图片将显示为占位符')
        except Exception as e:
            logger.warning(f'提取PPT图片失败: {str(e)}')
        
        return self.image_map
    
    def _extract_image(self, zf, name: str):
        """从zip中提取并转换图片"""
        import base64
        from PIL import Image
        import io
        
        try:
            img_data = zf.read(name)
            ext = name.split('.')[-1].lower()
            
            img_format = ext.upper()
            if img_format == 'JPG':
                img_format = 'JPEG'
            
            try:
                img = Image.open(io.BytesIO(img_data))
                
                if img.mode == 'RGBA' or img.mode == 'P':
                    pass
                elif img.mode != 'RGB':
                    img = img.convert('RGB')
                
                if img.width > self.max_image_size[0] or img.height > self.max_image_size[1]:
                    img.thumbnail(self.max_image_size, Image.Resampling.LANCZOS)
                
                output = io.BytesIO()
                img.save(output, format=img_format if img_format != 'JPEG' else 'JPEG', 
                        quality=self.image_quality, optimize=True)
                
                b64_data = base64.b64encode(output.getvalue()).decode('utf-8')
                mime_type = f'image/{ext.lower()}' if ext in ['png', 'jpeg', 'jpg', 'gif', 'bmp', 'webp'] else 'image/png'
                
                self.image_map[name] = {
                    'data': f'data:{mime_type};base64,{b64_data}',
                    'width': img.width,
                    'height': img.height,
                    'format': img_format
                }
            except ImportError:
                logger.warning(f'Pillow未安装，无法处理图片 {name}')
            except Exception as img_error:
                logger.warning(f'处理图片格式失败 {name}: {str(img_error)}')
                
        except Exception as e:
             logger.warning(f'提取图片失败 {name}: {str(e)}')


class PPTStyleExtractor:
    """PPT样式提取器 - 提取字体、背景、形状样式"""
    
    @staticmethod
    def rgb_to_hex(rgb_color) -> str:
        """将RGB颜色转换为十六进制"""
        if rgb_color is None:
            return None
        try:
            if hasattr(rgb_color, 'rgb'):
                r, g, b = rgb_color.rgb
                return '#{:02x}{:02x}{:02x}'.format(r, g, b)
            elif hasattr(rgb_color, 'theme') and rgb_color.theme:
                logger.debug('主题颜色，跳过')
                return None
        except Exception:
            pass
        return None
    
    @staticmethod
    def extract_font_style(font) -> dict:
        """提取字体样式"""
        if not font:
            return {}
        return {
            'name': font.name if hasattr(font, 'name') and font.name else None,
            'size': font.size.pt if font.size and hasattr(font.size, 'pt') else None,
            'color': PPTStyleExtractor.rgb_to_hex(font.color) if hasattr(font, 'color') and font.color else None,
            'bold': font.bold if hasattr(font, 'bold') else False,
            'italic': font.italic if hasattr(font, 'italic') else False,
            'underline': font.underline if hasattr(font, 'underline') and font.underline else False
        }
    
    @staticmethod
    def extract_paragraph_style(para) -> dict:
        """提取段落样式"""
        if not para:
            return {}
        return {
            'alignment': str(para.alignment) if para.alignment else None,
            'space_before': para.space_before.pt if para.space_before and hasattr(para.space_before, 'pt') else 0,
            'space_after': para.space_after.pt if para.space_after and hasattr(para.space_after, 'pt') else 0,
            'level': para.level if hasattr(para, 'level') else 0
        }
    
    @staticmethod
    def extract_shape_style(shape) -> dict:
        """提取形状样式"""
        style = {
            'fill': None,
            'line': None,
            'shadow': None
        }
        
        try:
            if hasattr(shape, 'fill') and shape.fill:
                fill = shape.fill
                fill_style = {'type': str(fill.type)}
                if fill.type != 1:
                    fill_style['color'] = PPTStyleExtractor.rgb_to_hex(fill.fore_color)
                style['fill'] = fill_style
            
            if hasattr(shape, 'line') and shape.line:
                line = shape.line
                style['line'] = {
                    'color': PPTStyleExtractor.rgb_to_hex(line.color),
                    'width': line.width.pt if line.width and hasattr(line.width, 'pt') else 0
                }
            
            if hasattr(shape, 'shadow') and shape.shadow:
                shadow = shape.shadow
                if hasattr(shadow, 'visible') and shadow.visible:
                    style['shadow'] = {
                        'type': str(shadow.type),
                        'visible': True
                    }
        except Exception as e:
            logger.debug(f'提取形状样式失败: {str(e)}')
        
        return style
    
    @staticmethod
    def extract_background(slide) -> dict:
        """提取幻灯片背景"""
        background = {
            'type': 'none',
            'color': None,
            'gradient': None
        }
        
        try:
            if hasattr(slide, 'background') and slide.background:
                bg = slide.background
                if hasattr(bg, 'fill') and bg.fill:
                    fill = bg.fill
                    if fill.type == 1:
                        background['type'] = 'solid'
                        background['color'] = PPTStyleExtractor.rgb_to_hex(fill.fore_color)
        except Exception as e:
            logger.debug(f'提取背景失败: {str(e)}')
        
        return background
    
    @staticmethod
    def extract_position(shape) -> dict:
        """提取形状位置和尺寸"""
        return {
            'left': shape.left.pt if hasattr(shape, 'left') and shape.left else 0,
            'top': shape.top.pt if hasattr(shape, 'top') and shape.top else 0,
            'width': shape.width.pt if hasattr(shape, 'width') and shape.width else 0,
            'height': shape.height.pt if hasattr(shape, 'height') and shape.height else 0
        }


def timeout_handler(seconds=30):
    """超时处理装饰器"""
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            result = [None]
            exception = [None]
            
            def target():
                try:
                    result[0] = func(*args, **kwargs)
                except Exception as e:
                    exception[0] = e
            
            thread = threading.Thread(target=target)
            thread.daemon = True
            thread.start()
            thread.join(seconds)
            
            if thread.is_alive():
                raise TimeoutError(f'操作超时（{seconds}秒）')
            
            if exception[0]:
                raise exception[0]
            
            return result[0]
        return wrapper
    return decorator


class TimeoutError(Exception):
    """预览处理超时异常"""
    pass


class OfficePreviewHandler:
    """Office文档预览处理器 - 支持多种转换格式和版本兼容性"""
    
    # 支持的Office文件扩展名
    SUPPORTED_EXTENSIONS = {
        'word': ['.doc', '.docx'],
        'excel': ['.xls', '.xlsx'],
        'powerpoint': ['.ppt', '.pptx']
    }
    
    # 最大文件大小限制（字节）
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    
    # 转换选项
    CONVERSION_OPTIONS = {
        'html': True,
        'pdf': True,
        'image': True,
        'text': True
    }
    
    @classmethod
    def preview_office_file(cls, file_path: str, conversion_format: str = None) -> dict:
        """预览Office文档 - 支持多种转换格式和版本兼容性
        
        Args:
            file_path: 文件路径
            conversion_format: 转换格式 (html, pdf, text, image)
        """
        try:
            file_name = os.path.basename(file_path)
            logger.info(f"开始预览Office文件: {file_path}, 文件名: {file_name}")
            
            # 1. 基础验证
            validation_result = cls._validate_file(file_path, file_name)
            if validation_result:
                return validation_result
            
            # 2. 确定文件类型
            file_ext = os.path.splitext(file_path)[1].lower()
            office_type = cls._get_office_type(file_ext)
            
            if not office_type:
                return cls._create_error_response(file_name, f'不支持的Office文件类型: {file_ext}')
            
            logger.info(f"文件类型: {office_type}, 扩展名: {file_ext}")
            
            # 如果指定了转换格式，执行转换
            if conversion_format:
                return cls._convert_office_file(file_path, file_name, office_type, file_ext, conversion_format)
            
            # 3. 根据文件类型选择处理方法
            preview_result = cls._process_office_file(file_path, file_name, office_type, file_ext)
            
            # 4. 添加通用信息
            preview_result.update({
                'file_size': os.path.getsize(file_path),
                'file_extension': file_ext,
                'conversion_supported': cls.CONVERSION_OPTIONS,
                'timestamp': cls._get_timestamp()
            })
            
            logger.info(f"Office文件预览完成: {file_name}")
            return preview_result
            
        except Exception as e:
            logger.error(f'Office文件预览失败: {str(e)}', exc_info=True)
            return cls._create_error_response(
                os.path.basename(file_path) if file_path else '未知文件',
                f'预览失败: {str(e)}'
            )
    
    @classmethod
    def _convert_office_file(cls, file_path: str, file_name: str, office_type: str, file_ext: str, conversion_format: str) -> dict:
        """转换Office文档到指定格式"""
        try:
            # 检查转换格式是否支持
            supported_formats = cls.CONVERSION_OPTIONS
            if conversion_format not in supported_formats or not supported_formats[conversion_format]:
                return cls._create_error_response(file_name, f'不支持将{office_type}文档转换为{conversion_format}格式')
            
            # 执行转换
            if conversion_format == 'html':
                return cls._convert_to_html(file_path, file_name, office_type, file_ext)
            elif conversion_format == 'pdf':
                return cls._convert_to_pdf(file_path, file_name, office_type, file_ext)
            elif conversion_format == 'text':
                return cls._convert_to_text(file_path, file_name, office_type, file_ext)
            elif conversion_format == 'image':
                return cls._convert_to_image(file_path, file_name, office_type, file_ext)
            else:
                return cls._create_error_response(file_name, f'不支持的转换格式: {conversion_format}')
                
        except Exception as e:
            logger.error(f'Office文档转换失败: {str(e)}', exc_info=True)
            return cls._create_error_response(file_name, f'Office文档转换失败: {str(e)}')
    
    @classmethod
    def _validate_file(cls, file_path: str, file_name: str) -> Optional[dict]:
        """验证文件基础信息"""
        # 检查文件是否存在
        if not os.path.exists(file_path):
            logger.error(f"文件不存在: {file_path}")
            return cls._create_error_response(file_name, '文件不存在')
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        if file_size > cls.MAX_FILE_SIZE:
            logger.warning(f"文件过大: {file_size}字节")
            return cls._create_error_response(
                file_name, 
                f'文件过大（{cls._format_file_size(file_size)}），请下载后查看'
            )
        
        # 检查文件是否可读
        if not os.access(file_path, os.R_OK):
            logger.error(f"文件不可读: {file_path}")
            return cls._create_error_response(file_name, '文件不可读')
        
        return None
    
    @classmethod
    def _get_office_type(cls, file_ext: str) -> Optional[str]:
        """根据文件扩展名获取Office类型"""
        for office_type, extensions in cls.SUPPORTED_EXTENSIONS.items():
            if file_ext in extensions:
                return office_type
        return None
    
    @classmethod
    def _process_office_file(cls, file_path: str, file_name: str, office_type: str, file_ext: str) -> dict:
        """处理Office文件预览"""
        try:
            # 根据文件类型调用相应的处理方法
            if office_type == 'word':
                return cls._preview_word_file(file_path, file_name, file_ext)
            elif office_type == 'excel':
                return cls._preview_excel_file(file_path, file_name, file_ext)
            elif office_type == 'powerpoint':
                return cls._preview_powerpoint_file(file_path, file_name, file_ext)
            else:
                raise ValueError(f'不支持的Office类型: {office_type}')
                
        except Exception as e:
            logger.error(f'Office文件处理失败: {str(e)}', exc_info=True)
            return cls._create_error_response(file_name, f'Office文件处理失败: {str(e)}')
    
    @classmethod
    def _create_error_response(cls, file_name: str, message: str) -> dict:
        """创建错误响应"""
        return {
            'type': 'office_enhanced',
            'name': file_name,
            'office_type': 'error',
            'preview_options': [
                {
                    'type': 'text',
                    'content': message,
                    'name': '错误信息'
                }
            ]
        }
    
    @classmethod
    def _convert_to_html(cls, file_path: str, file_name: str, office_type: str, file_ext: str) -> dict:
        """将Office文档转换为HTML格式"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                if office_type == 'word':
                    return cls._convert_word_to_html(file_path, file_name, file_ext, temp_dir)
                elif office_type == 'excel':
                    return cls._convert_excel_to_html(file_path, file_name, file_ext, temp_dir)
                elif office_type == 'powerpoint':
                    return cls._convert_powerpoint_to_html(file_path, file_name, file_ext, temp_dir)
                else:
                    raise ValueError(f'不支持的Office类型: {office_type}')
                    
        except Exception as e:
            logger.error(f'HTML转换失败: {str(e)}', exc_info=True)
            return cls._create_error_response(file_name, f'HTML转换失败: {str(e)}')
    
    @classmethod
    def _convert_word_to_html(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将Word文档转换为HTML"""
        try:
            if file_ext == '.docx':
                from docx import Document
                
                doc = Document(file_path)
                html_content = []
                
                # 处理段落
                for para in doc.paragraphs:
                    if para.text.strip():
                        html_content.append(f'<p>{para.text}</p>')
                
                # 处理表格
                for table in doc.tables:
                    html_content.append('<table border="1" style="border-collapse: collapse; width: 100%;">')
                    for row in table.rows:
                        html_content.append('<tr>')
                        for cell in row.cells:
                            html_content.append(f'<td style="padding: 8px;">{cell.text}</td>')
                        html_content.append('</tr>')
                    html_content.append('</table>')
                
                if not html_content:
                    html_content = ['<p>文档为空</p>']
                
                return {
                    'type': 'html',
                    'name': file_name,
                    'content': '\n'.join(html_content),
                    'file_format': 'html'
                }
            else:
                # 对于.doc文件，使用PyMuPDF提取文本
                import fitz
                
                doc = fitz.open(file_path)
                html_content = []
                
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    text = page.get_text()
                    
                    if text.strip():
                        # 将文本按段落分割
                        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
                        for para in paragraphs:
                            html_content.append(f'<p>{para}</p>')
                
                doc.close()
                
                if not html_content:
                    html_content = ['<p>文档为空</p>']
                
                return {
                    'type': 'html',
                    'name': file_name,
                    'content': '\n'.join(html_content),
                    'file_format': 'html',
                    'note': '旧版Word文档，格式可能不完整'
                }
                
        except Exception as e:
            raise Exception(f'Word文档HTML转换失败: {str(e)}')
    
    @classmethod
    def _convert_excel_to_html(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将Excel文档转换为HTML"""
        try:
            html_content = []
            
            if file_ext == '.xlsx':
                from openpyxl import load_workbook
                
                wb = load_workbook(file_path, data_only=True)
                
                for sheet_name in wb.sheetnames[:3]:  # 最多处理3个工作表
                    ws = wb[sheet_name]
                    
                    html_content.append(f'<h3>{sheet_name}</h3>')
                    html_content.append('<table border="1" style="border-collapse: collapse; width: 100%;">')
                    
                    # 获取数据范围
                    max_row = min(20, ws.max_row)  # 最多20行
                    max_col = min(10, ws.max_column)  # 最多10列
                    
                    for row in range(1, max_row + 1):
                        html_content.append('<tr>')
                        for col in range(1, max_col + 1):
                            cell_value = ws.cell(row=row, column=col).value
                            if cell_value is None:
                                cell_value = ''
                            html_content.append(f'<td style="padding: 8px;">{cell_value}</td>')
                        html_content.append('</tr>')
                    html_content.append('</table>')
                
                wb.close()
                
            else:  # .xls
                import xlrd
                
                wb = xlrd.open_workbook(file_path)
                
                for sheet_index in range(min(3, wb.nsheets)):  # 最多处理3个工作表
                    ws = wb.sheet_by_index(sheet_index)
                    
                    html_content.append(f'<h3>{ws.name}</h3>')
                    html_content.append('<table border="1" style="border-collapse: collapse; width: 100%;">')
                    
                    # 获取数据范围
                    max_row = min(20, ws.nrows)  # 最多20行
                    max_col = min(10, ws.ncols)  # 最多10列
                    
                    for row in range(max_row):
                        html_content.append('<tr>')
                        for col in range(max_col):
                            cell_value = ws.cell_value(row, col)
                            if cell_value is None:
                                cell_value = ''
                            html_content.append(f'<td style="padding: 8px;">{cell_value}</td>')
                        html_content.append('</tr>')
                    html_content.append('</table>')
                
                wb.release_resources()
            
            if not html_content:
                html_content = ['<p>工作簿为空</p>']
            
            return {
                'type': 'html',
                'name': file_name,
                'content': '\n'.join(html_content),
                'file_format': 'html'
            }
                
        except Exception as e:
            raise Exception(f'Excel文档HTML转换失败: {str(e)}')
    
    @classmethod
    def _convert_powerpoint_to_html(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将PowerPoint文档转换为HTML"""
        try:
            html_content = []
            
            if file_ext == '.pptx':
                from pptx import Presentation
                
                prs = Presentation(file_path)
                
                for i, slide in enumerate(prs.slides[:10]):  # 最多10张幻灯片
                    html_content.append(f'<h3>幻灯片 {i+1}</h3>')
                    html_content.append('<div style="border: 1px solid #ccc; padding: 20px; margin: 10px 0;">')
                    
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        for text in slide_text:
                            html_content.append(f'<p>{text}</p>')
                    else:
                        html_content.append('<p>此幻灯片无文本内容</p>')
                    
                    html_content.append('</div>')
                
            else:  # .ppt
                import fitz
                
                doc = fitz.open(file_path)
                
                for page_num in range(min(10, len(doc))):
                    page = doc.load_page(page_num)
                    text = page.get_text()
                    
                    html_content.append(f'<h3>幻灯片 {page_num+1}</h3>')
                    html_content.append('<div style="border: 1px solid #ccc; padding: 20px; margin: 10px 0;">')
                    
                    if text.strip():
                        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
                        for para in paragraphs:
                            html_content.append(f'<p>{para}</p>')
                    else:
                        html_content.append('<p>此幻灯片无文本内容</p>')
                    
                    html_content.append('</div>')
                
                doc.close()
            
            if not html_content:
                html_content = ['<p>演示文稿为空</p>']
            
            return {
                'type': 'html',
                'name': file_name,
                'content': '\n'.join(html_content),
                'file_format': 'html'
            }
                
        except Exception as e:
            raise Exception(f'PowerPoint文档HTML转换失败: {str(e)}')
    
    @classmethod
    def _format_file_size(cls, size_bytes: int) -> str:
        """格式化文件大小"""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024.0:
                return f"{size_bytes:.2f} {unit}"
            size_bytes /= 1024.0
        return f"{size_bytes:.2f} TB"
    
    @classmethod
    def _convert_to_pdf(cls, file_path: str, file_name: str, office_type: str, file_ext: str) -> dict:
        """将Office文档转换为PDF格式"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                if office_type == 'word':
                    return cls._convert_word_to_pdf(file_path, file_name, file_ext, temp_dir)
                elif office_type == 'excel':
                    return cls._convert_excel_to_pdf(file_path, file_name, file_ext, temp_dir)
                elif office_type == 'powerpoint':
                    return cls._convert_powerpoint_to_pdf(file_path, file_name, file_ext, temp_dir)
                else:
                    raise ValueError(f'不支持的Office类型: {office_type}')
                    
        except Exception as e:
            logger.error(f'PDF转换失败: {str(e)}', exc_info=True)
            return cls._create_error_response(file_name, f'PDF转换失败: {str(e)}')
    
    @classmethod
    def _convert_word_to_pdf(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将Word文档转换为PDF"""
        try:
            # 对于Word文档，使用PyMuPDF进行转换
            import fitz
            
            doc = fitz.open(file_path)
            
            # 创建PDF文件
            pdf_path = os.path.join(temp_dir, f"{os.path.splitext(file_name)[0]}.pdf")
            
            # 使用PyMuPDF直接保存为PDF
            doc.save(pdf_path)
            doc.close()
            
            # 读取PDF内容
            with open(pdf_path, 'rb') as f:
                pdf_content = f.read()
            
            return {
                'type': 'pdf',
                'name': f"{os.path.splitext(file_name)[0]}.pdf",
                'content': pdf_content,
                'file_format': 'pdf',
                'size': len(pdf_content)
            }
                
        except Exception as e:
            raise Exception(f'Word文档PDF转换失败: {str(e)}')
    
    @classmethod
    def _convert_excel_to_pdf(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将Excel文档转换为PDF"""
        try:
            # 对于Excel文档，先转换为HTML，然后使用PyMuPDF生成PDF
            html_result = cls._convert_excel_to_html(file_path, file_name, file_ext, temp_dir)
            
            if html_result['type'] == 'html':
                # 使用PyMuPDF从HTML生成PDF
                import fitz
                
                # 创建PDF文档
                pdf_doc = fitz.open()
                
                # 添加HTML内容
                page = pdf_doc.new_page()
                
                # 设置页面尺寸
                page.set_mediabox(fitz.Rect(0, 0, 595, 842))  # A4尺寸
                
                # 插入HTML内容
                page.insert_html(html_result['content'])
                
                # 保存PDF
                pdf_path = os.path.join(temp_dir, f"{os.path.splitext(file_name)[0]}.pdf")
                pdf_doc.save(pdf_path)
                pdf_doc.close()
                
                # 读取PDF内容
                with open(pdf_path, 'rb') as f:
                    pdf_content = f.read()
                
                return {
                    'type': 'pdf',
                    'name': f"{os.path.splitext(file_name)[0]}.pdf",
                    'content': pdf_content,
                    'file_format': 'pdf',
                    'size': len(pdf_content),
                    'note': 'Excel文档转换为PDF，格式可能有所调整'
                }
            else:
                raise Exception('HTML转换失败，无法生成PDF')
                
        except Exception as e:
            raise Exception(f'Excel文档PDF转换失败: {str(e)}')
    
    @classmethod
    def _convert_powerpoint_to_pdf(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将PowerPoint文档转换为PDF"""
        try:
            # 对于PowerPoint文档，使用PyMuPDF进行转换
            import fitz
            
            doc = fitz.open(file_path)
            
            # 创建PDF文件
            pdf_path = os.path.join(temp_dir, f"{os.path.splitext(file_name)[0]}.pdf")
            
            # 使用PyMuPDF直接保存为PDF
            doc.save(pdf_path)
            doc.close()
            
            # 读取PDF内容
            with open(pdf_path, 'rb') as f:
                pdf_content = f.read()
            
            return {
                'type': 'pdf',
                'name': f"{os.path.splitext(file_name)[0]}.pdf",
                'content': pdf_content,
                'file_format': 'pdf',
                'size': len(pdf_content)
            }
                
        except Exception as e:
            raise Exception(f'PowerPoint文档PDF转换失败: {str(e)}')
    
    @classmethod
    def _convert_to_text(cls, file_path: str, file_name: str, office_type: str, file_ext: str) -> dict:
        """将Office文档转换为纯文本格式"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                if office_type == 'word':
                    return cls._convert_word_to_text(file_path, file_name, file_ext, temp_dir)
                elif office_type == 'excel':
                    return cls._convert_excel_to_text(file_path, file_name, file_ext, temp_dir)
                elif office_type == 'powerpoint':
                    return cls._convert_powerpoint_to_text(file_path, file_name, file_ext, temp_dir)
                else:
                    raise ValueError(f'不支持的Office类型: {office_type}')
                    
        except Exception as e:
            logger.error(f'文本转换失败: {str(e)}', exc_info=True)
            return cls._create_error_response(file_name, f'文本转换失败: {str(e)}')
    
    @classmethod
    def _convert_word_to_text(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将Word文档转换为纯文本"""
        try:
            text_content = []
            
            if file_ext == '.docx':
                from docx import Document
                
                doc = Document(file_path)
                
                # 处理段落
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_content.append(para.text)
                
                # 处理表格
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text)
                        text_content.append('\t'.join(row_text))
                
            else:  # .doc
                import fitz
                
                doc = fitz.open(file_path)
                
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    text = page.get_text()
                    
                    if text.strip():
                        text_content.append(text)
                
                doc.close()
            
            if not text_content:
                text_content = ['文档为空']
            
            return {
                'type': 'text',
                'name': f"{os.path.splitext(file_name)[0]}.txt",
                'content': '\n\n'.join(text_content),
                'file_format': 'txt'
            }
                
        except Exception as e:
            raise Exception(f'Word文档文本转换失败: {str(e)}')
    
    @classmethod
    def _convert_excel_to_text(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将Excel文档转换为纯文本"""
        try:
            text_content = []
            
            if file_ext == '.xlsx':
                from openpyxl import load_workbook
                
                wb = load_workbook(file_path, data_only=True)
                
                for sheet_name in wb.sheetnames[:3]:  # 最多处理3个工作表
                    ws = wb[sheet_name]
                    
                    text_content.append(f'=== {sheet_name} ===')
                    
                    # 获取数据范围
                    max_row = min(20, ws.max_row)  # 最多20行
                    max_col = min(10, ws.max_column)  # 最多10列
                    
                    for row in range(1, max_row + 1):
                        row_data = []
                        for col in range(1, max_col + 1):
                            cell_value = ws.cell(row=row, column=col).value
                            if cell_value is None:
                                cell_value = ''
                            row_data.append(str(cell_value))
                        text_content.append('\t'.join(row_data))
                    
                    text_content.append('')  # 空行分隔
                
                wb.close()
                
            else:  # .xls
                import xlrd
                
                wb = xlrd.open_workbook(file_path)
                
                for sheet_index in range(min(3, wb.nsheets)):  # 最多处理3个工作表
                    ws = wb.sheet_by_index(sheet_index)
                    
                    text_content.append(f'=== {ws.name} ===')
                    
                    # 获取数据范围
                    max_row = min(20, ws.nrows)  # 最多20行
                    max_col = min(10, ws.ncols)  # 最多10列
                    
                    for row in range(max_row):
                        row_data = []
                        for col in range(max_col):
                            cell_value = ws.cell_value(row, col)
                            if cell_value is None:
                                cell_value = ''
                            row_data.append(str(cell_value))
                        text_content.append('\t'.join(row_data))
                    
                    text_content.append('')  # 空行分隔
                
                wb.release_resources()
            
            if not text_content:
                text_content = ['工作簿为空']
            
            return {
                'type': 'text',
                'name': f"{os.path.splitext(file_name)[0]}.txt",
                'content': '\n'.join(text_content),
                'file_format': 'txt'
            }
                
        except Exception as e:
            raise Exception(f'Excel文档文本转换失败: {str(e)}')
    
    @classmethod
    def _convert_powerpoint_to_text(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将PowerPoint文档转换为纯文本"""
        try:
            text_content = []
            
            if file_ext == '.pptx':
                from pptx import Presentation
                
                prs = Presentation(file_path)
                
                for i, slide in enumerate(prs.slides[:10]):  # 最多10张幻灯片
                    text_content.append(f'=== 幻灯片 {i+1} ===')
                    
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        text_content.extend(slide_text)
                    else:
                        text_content.append('此幻灯片无文本内容')
                    
                    text_content.append('')  # 空行分隔
                
            else:  # .ppt
                import fitz
                
                doc = fitz.open(file_path)
                
                for page_num in range(min(10, len(doc))):
                    page = doc.load_page(page_num)
                    text = page.get_text()
                    
                    text_content.append(f'=== 幻灯片 {page_num+1} ===')
                    
                    if text.strip():
                        text_content.append(text)
                    else:
                        text_content.append('此幻灯片无文本内容')
                    
                    text_content.append('')  # 空行分隔
                
                doc.close()
            
            if not text_content:
                text_content = ['演示文稿为空']
            
            return {
                'type': 'text',
                'name': f"{os.path.splitext(file_name)[0]}.txt",
                'content': '\n'.join(text_content),
                'file_format': 'txt'
            }
                
        except Exception as e:
            raise Exception(f'PowerPoint文档文本转换失败: {str(e)}')
    
    @classmethod
    def _get_timestamp(cls) -> str:
        """获取当前时间戳"""
        import datetime
        return datetime.datetime.now().isoformat()
    
    @classmethod
    def _preview_word_file(cls, file_path: str, file_name: str, file_ext: str) -> dict:
        """预览Word文档 - 支持.doc和.docx格式"""
        try:
            # 根据文件格式选择处理方法
            if file_ext == '.doc':
                return cls._preview_old_word_file(file_path, file_name)
            else:  # .docx
                return cls._preview_new_word_file(file_path, file_name)
                
        except Exception as e:
            logger.error(f'Word文档预览失败: {str(e)}', exc_info=True)
            return cls._create_error_response(file_name, f'Word文档预览失败: {str(e)}')
    
    @classmethod
    def _preview_new_word_file(cls, file_path: str, file_name: str) -> dict:
        """预览新版Word文档(.docx) - 保留格式信息和图片"""
        try:
            from docx import Document
            from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
            from docx.oxml.ns import qn
            import base64
            import uuid
            import os
            from django.conf import settings
            
            file_size = os.path.getsize(file_path)
            if file_size > 20 * 1024 * 1024:
                logger.warning(f"文件过大: {file_size} bytes")
                return {
                    'type': 'office_enhanced',
                    'office_type': 'word',
                    'name': file_name,
                    'preview_options': [
                        {
                            'type': 'text',
                            'content': f'文件过大（{file_size}字节），请下载后查看',
                            'name': '提示信息'
                        }
                    ],
                    'conversion_available': ['pdf', 'text']
                }
            
            doc = Document(file_path)
            
            preview_options = []
            content_blocks = []
            images = []
            
            for para in doc.paragraphs:
                text = TextCleaner.clean_text(para.text)
                if not text:
                    continue
                
                style_name = para.style.name if para.style else 'Normal'
                alignment = para.alignment
                is_heading = False
                heading_level = 0
                
                if style_name.startswith('Heading'):
                    is_heading = True
                    try:
                        heading_level = int(style_name.replace('Heading', ''))
                    except:
                        heading_level = 1
                
                para_info = {
                    'text': text,
                    'style': style_name,
                    'is_heading': is_heading,
                    'level': heading_level,
                    'alignment': str(alignment) if alignment else None,
                    'runs': []
                }
                
                for run in para.runs:
                    run_text = TextCleaner.clean_run_text(run.text)
                    if not run_text:
                        continue
                    
                    run_info = {
                        'text': run_text,
                        'bold': run.bold if hasattr(run, 'bold') else False,
                        'italic': run.italic if hasattr(run, 'italic') else False,
                        'underline': run.underline if hasattr(run, 'underline') else False,
                        'font_size': None
                    }
                    
                    if run.font and hasattr(run.font, 'size') and run.font.size:
                        run_info['font_size'] = run.font.size.pt
                    
                    para_info['runs'].append(run_info)
                
                if para_info['runs']:
                    content_blocks.append(para_info)
            
            if content_blocks:
                preview_options.append({
                    'type': 'formatted_text',
                    'content': content_blocks,
                    'name': '文档内容',
                    'format': 'structured'
                })
            
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                has_content = False
                
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = TextCleaner.clean_text(cell.text)
                        row_data.append(cell_text)
                        if cell_text:
                            has_content = True
                    table_data.append(row_data)
                
                if has_content and table_data:
                    headers = table_data[0] if len(table_data) > 0 else []
                    data = table_data[1:] if len(table_data) > 1 else []
                    
                    preview_options.append({
                        'type': 'table',
                        'content': {
                            'headers': headers,
                            'data': data
                        },
                        'name': f'表格 {table_idx + 1}',
                        'table_index': table_idx + 1
                    })
            
            if doc.inline_shapes:
                for idx, shape in enumerate(doc.inline_shapes):
                    try:
                        image_info = {
                            'index': idx + 1,
                            'type': 'image',
                            'width': None,
                            'height': None,
                            'alt_text': ''
                        }
                        
                        if hasattr(shape, 'width') and shape.width:
                            image_info['width'] = shape.width.pt
                        if hasattr(shape, 'height') and shape.height:
                            image_info['height'] = shape.height.pt
                        
                        if hasattr(shape, 'alt_text') and shape.alt_text:
                            image_info['alt_text'] = shape.alt_text
                        
                        images.append(image_info)
                    except Exception as img_error:
                        logger.warning(f'处理Word图片 {idx+1} 失败: {str(img_error)}')
                        continue
            
            if images:
                preview_options.append({
                    'type': 'images',
                    'content': images,
                    'name': f'图片 ({len(images)}张)',
                    'format': 'gallery'
                })
            
            if not preview_options:
                preview_options.append({
                    'type': 'text',
                    'content': '文档为空或无法提取内容',
                    'name': '提示信息'
                })
            
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'word',
                'file_format': 'docx',
                'preview_options': preview_options,
                'conversion_available': ['html', 'pdf', 'text']
            }
            
        except Exception as e:
            logger.error(f'新版Word文档处理失败: {str(e)}', exc_info=True)
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'word',
                'file_format': 'docx',
                'preview_options': [{
                    'type': 'text',
                    'content': 'Word文档内容提取失败，建议下载后查看原文件。错误详情：' + str(e),
                    'name': '错误信息'
                }],
                'conversion_available': ['pdf', 'text'],
                'error': str(e)
            }
    
    @classmethod
    def _preview_old_word_file(cls, file_path: str, file_name: str) -> dict:
        """预览旧版Word文档(.doc)"""
        try:
            import fitz
            
            file_size = os.path.getsize(file_path)
            if file_size > 20 * 1024 * 1024:
                logger.warning(f"文件过大: {file_size} bytes")
                return {
                    'type': 'office_enhanced',
                    'office_type': 'word',
                    'name': file_name,
                    'preview_options': [
                        {
                            'type': 'text',
                            'content': f'文件过大（{file_size}字节），请下载后查看',
                            'name': '提示信息'
                        }
                    ],
                    'conversion_available': ['text'],
                    'compatibility_note': '旧版Word文档'
                }
            
            doc = fitz.open(file_path)
            text_content = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    text_content.append(text)
            
            doc.close()
            
            preview_options = []
            
            if text_content:
                preview_options.append({
                    'type': 'text',
                    'content': '\n\n'.join(text_content),
                    'name': '文档内容',
                    'format': 'plain'
                })
            else:
                preview_options.append({
                    'type': 'text',
                    'content': '无法提取文档内容，请下载后查看',
                    'name': '提示信息'
                })
            
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'word',
                'file_format': 'doc',
                'preview_options': preview_options,
                'conversion_available': ['text'],
                'compatibility_note': '旧版Word文档，格式可能不完整'
            }
            
        except Exception as e:
            logger.error(f'旧版Word文档处理失败: {str(e)}', exc_info=True)
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'word',
                'file_format': 'doc',
                'preview_options': [{
                    'type': 'text',
                    'content': '旧版Word文档预览失败，建议下载后查看。错误详情：' + str(e),
                    'name': '错误信息'
                }],
                'conversion_available': ['text'],
                'error': str(e)
            }
    
    @classmethod
    def _preview_excel_file(cls, file_path: str, file_name: str, file_ext: str) -> dict:
        """预览Excel文档 - 支持.xls和.xlsx格式"""
        try:
            # 首先尝试使用PyMuPDF作为备选方法，如果其他方法失败
            try:
                import fitz
                doc = fitz.open(file_path)
                if doc.page_count > 0:
                    # 使用PyMuPDF提取文本内容作为备选
                    text_content = []
                    for page_num in range(min(5, doc.page_count)):  # 最多5页
                        page = doc.load_page(page_num)
                        text = page.get_text()
                        if text.strip():
                            text_content.append(f'=== 页 {page_num+1} ===\n{text}')
                    doc.close()
                    
                    if text_content:
                        # 返回PyMuPDF提取的结果
                        return {
                            'type': 'office_enhanced',
                            'name': file_name,
                            'office_type': 'excel',
                            'file_format': file_ext,
                            'preview_options': [{
                                'type': 'text',
                                'content': '\n'.join(text_content) + '\n\n[提示: 使用替代方法提取的内容]',
                                'name': '表格内容预览',
                                'format': 'plain'
                            }],
                            'conversion_available': ['html', 'csv', 'json', 'text']
                        }
            except Exception as fitz_error:
                logger.warning(f'PyMuPDF处理Excel失败: {str(fitz_error)}')
            
            # 根据文件格式选择标准处理方法
            if file_ext == '.xls':
                return cls._preview_old_excel_file(file_path, file_name)
            else:  # .xlsx
                return cls._preview_new_excel_file(file_path, file_name)
                
        except Exception as e:
            logger.error(f'Excel文档预览失败: {str(e)}', exc_info=True)
            try:
                with tempfile.TemporaryDirectory() as temp_dir:
                    text_result = cls._convert_excel_to_text(file_path, file_name, file_ext, temp_dir)
                    if text_result and 'content' in text_result:
                        return {
                            'type': 'office_enhanced',
                            'name': file_name,
                            'office_type': 'excel',
                            'file_format': file_ext,
                            'preview_options': [{
                                'type': 'text',
                                'content': text_result['content'] + '\n\n[注意: 由于文档格式问题，仅显示文本内容]',
                                'name': '表格文本内容',
                                'format': 'plain'
                            }],
                            'conversion_available': ['text'],
                            'error': str(e)
                        }
            except Exception as alt_e:
                logger.warning(f'备选Excel预览方案失败: {str(alt_e)}')
            
            return cls._create_error_response(file_name, f'Excel文档预览失败: {str(e)}')
    
    @classmethod
    def _preview_new_excel_file(cls, file_path: str, file_name: str) -> dict:
        """预览新版Excel文档(.xlsx)"""
        try:
            from openpyxl import load_workbook
            
            # 检查文件大小
            file_size = os.path.getsize(file_path)
            if file_size > 20 * 1024 * 1024:  # 20MB
                logger.warning(f"Excel文件过大: {file_size} bytes")
                
            wb = load_workbook(file_path, read_only=True, data_only=True)
            
            # 获取工作表信息
            sheet_names = wb.sheetnames
            
            # 预览前3个工作表，每个工作表最多10行数据
            preview_options = []
            
            # 为每个工作表添加单独的异常处理
            for i, sheet_name in enumerate(sheet_names[:3]):  # 最多预览3个工作表
                try:
                    ws = wb[sheet_name]
                    
                    # 获取表头和数据
                    headers = []
                    data = []
                    
                    # 读取前10行数据
                    for row_num, row in enumerate(ws.iter_rows(max_row=10, values_only=True)):
                        try:
                            if row_num == 0:
                                # 第一行作为表头
                                headers = [str(cell) if cell is not None else '' for cell in row]
                            else:
                                # 数据行
                                data_row = [str(cell) if cell is not None else '' for cell in row]
                                if any(data_row):  # 只添加非空行
                                    data.append(data_row)
                        except Exception as row_error:
                            logger.warning(f'工作表 {sheet_name} 的行 {row_num+1} 处理失败: {str(row_error)}')
                            continue
                    
                    # 构建表格预览
                    if headers or data:
                        preview_options.append({
                            'type': 'table',
                            'content': {
                                'headers': headers,
                                'data': data
                            },
                            'name': f'工作表: {sheet_name}',
                            'sheet_info': {
                                'name': sheet_name,
                                'index': i + 1,
                                'total_sheets': len(sheet_names)
                            },
                            'sheets': [{'name': sheet_name, 'hasHeader': True if headers else False, 'data': [headers] + data}]
                        })
                except Exception as sheet_error:
                    logger.warning(f'工作表 {sheet_name} 处理失败: {str(sheet_error)}')
                    continue
            
            # 如果工作表数量超过3个，添加提示
            if len(sheet_names) > 3:
                preview_options.append({
                    'type': 'text',
                    'content': f'文档包含{len(sheet_names)}个工作表，仅显示前3个',
                    'name': '提示信息'
                })
            
            # 如果没有成功预览任何工作表，尝试添加一个基本提示
            if not preview_options:
                preview_options.append({
                    'type': 'text',
                    'content': f'无法提取工作表内容，工作簿包含 {len(sheet_names)} 个工作表',
                    'name': '提示信息'
                })
            
            wb.close()
            
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'excel',
                'file_format': 'xlsx',
                'preview_options': preview_options,
                'conversion_available': ['html', 'csv', 'json']
            }
            
        except Exception as e:
            logger.error(f'新版Excel文档处理失败: {str(e)}', exc_info=True)
            # 即使处理失败也返回基本信息，而不是直接抛出异常
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'excel',
                'file_format': 'xlsx',
                'preview_options': [{
                    'type': 'text',
                    'content': 'Excel文档内容提取失败，建议下载后查看原文件。错误详情：' + str(e),
                    'name': '错误信息'
                }],
                'conversion_available': ['text'],
                'error': str(e)
            }
    
    @classmethod
    def _preview_old_excel_file(cls, file_path: str, file_name: str) -> dict:
        """预览旧版Excel文档(.xls)"""
        try:
            import xlrd
            
            # 打开工作簿，添加格式设置以处理不同版本
            wb = xlrd.open_workbook(file_path, formatting_info=False)
            
            # 获取工作表信息
            sheet_names = wb.sheet_names()
            
            # 预览前3个工作表，每个工作表最多10行数据
            preview_options = []
            
            # 为每个工作表添加单独的异常处理
            for i, sheet_name in enumerate(sheet_names[:3]):
                try:
                    ws = wb.sheet_by_index(i)
                    
                    # 获取表头和数据
                    headers = []
                    data = []
                    rows_data = []
                    
                    # 读取前10行数据，添加单独的行异常处理
                    for row_num in range(min(10, ws.nrows)):
                        try:
                            row_data = []
                            for col_num in range(min(10, ws.ncols)):  # 限制列数，避免处理过多列
                                try:
                                    cell_value = ws.cell_value(row_num, col_num)
                                    # 处理不同类型的单元格值
                                    if cell_value is None:
                                        cell_str = ''
                                    elif isinstance(cell_value, float) and cell_value.is_integer():
                                        cell_str = str(int(cell_value))
                                    else:
                                        cell_str = str(cell_value)
                                    row_data.append(cell_str)
                                except Exception as cell_error:
                                    logger.warning(f'工作表 {sheet_name} 的单元格 ({row_num},{col_num}) 处理失败: {str(cell_error)}')
                                    row_data.append('无法读取')
                                    continue
                            
                            rows_data.append(row_data)
                            
                            if row_num == 0:
                                headers = row_data
                            else:
                                if any(row_data):
                                    data.append(row_data)
                        except Exception as row_error:
                            logger.warning(f'工作表 {sheet_name} 的行 {row_num+1} 处理失败: {str(row_error)}')
                            continue
                    
                    # 构建表格预览
                    if headers or data:
                        preview_options.append({
                            'type': 'table',
                            'content': {
                                'headers': headers,
                                'data': data
                            },
                            'name': f'工作表: {sheet_name}',
                            'sheet_info': {
                                'name': sheet_name,
                                'index': i + 1,
                                'total_sheets': len(sheet_names)
                            },
                            'note': '旧版Excel文档，可能无法完全保留格式',
                            'sheets': [{'name': sheet_name, 'hasHeader': True if headers else False, 'data': rows_data}]
                        })
                except Exception as sheet_error:
                    logger.warning(f'工作表 {sheet_name} 处理失败: {str(sheet_error)}')
                    continue
            
            # 如果工作表数量超过3个，添加提示
            if len(sheet_names) > 3:
                preview_options.append({
                    'type': 'text',
                    'content': f'文档包含{len(sheet_names)}个工作表，仅显示前3个',
                    'name': '提示信息'
                })
            
            # 如果没有成功预览任何工作表，添加基本信息
            if not preview_options:
                preview_options.append({
                    'type': 'text',
                    'content': f'无法提取工作表内容，工作簿包含 {len(sheet_names)} 个工作表',
                    'name': '提示信息'
                })
            
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'excel',
                'file_format': 'xls',
                'preview_options': preview_options,
                'conversion_available': ['csv'],
                'compatibility_note': '旧版Excel文档，功能受限'
            }
            
        except Exception as e:
            logger.error(f'旧版Excel文档处理失败: {str(e)}', exc_info=True)
            # 即使处理失败也返回基本信息，而不是直接抛出异常
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'excel',
                'file_format': 'xls',
                'preview_options': [{
                    'type': 'text',
                    'content': '旧版Excel文档内容提取失败，建议下载后查看原文件。错误详情：' + str(e),
                    'name': '错误信息'
                }],
                'conversion_available': ['text'],
                'error': str(e)
            }
    
    @classmethod
    def _preview_powerpoint_file(cls, file_path: str, file_name: str, file_ext: str) -> dict:
        """预览PowerPoint文档 - 支持.ppt和.pptx格式"""
        try:
            # 根据文件格式选择处理方法
            if file_ext == '.ppt':
                return cls._preview_old_powerpoint_file(file_path, file_name)
            else:  # .pptx
                return cls._preview_new_powerpoint_file(file_path, file_name)
                
        except Exception as e:
            logger.error(f'PowerPoint文档预览失败: {str(e)}', exc_info=True)
            return cls._create_error_response(file_name, f'PowerPoint文档预览失败: {str(e)}')
    
    @classmethod
    def _preview_new_powerpoint_file(cls, file_path: str, file_name: str) -> dict:
        """预览新版PowerPoint文档(.pptx) - 保留格式和结构"""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE
            
            preview_options = []
            image_map = {}
            
            try:
                image_extractor = PPTImageExtractor(file_path)
                image_map = image_extractor.extract_all_images()
            except Exception as img_extract_error:
                logger.warning(f'PPT图片提取失败: {str(img_extract_error)}')
            
            try:
                prs = Presentation(file_path)
                slide_count = len(prs.slides)
                slides_to_preview = min(10, slide_count)
                
                for i in range(slides_to_preview):
                    try:
                        slide = prs.slides[i]
                    except Exception as slide_access_error:
                        logger.warning(f'访问幻灯片 {i+1} 失败: {str(slide_access_error)}')
                        break
                    
                    slide_content = {
                        'slide_number': i + 1,
                        'total_slides': slide_count,
                        'title': '',
                        'background': PPTStyleExtractor.extract_background(slide),
                        'content_blocks': [],
                        'shapes': []
                    }
                    
                    try:
                        if slide.shapes.title and slide.shapes.title.text.strip():
                            slide_content['title'] = TextCleaner.clean_text(slide.shapes.title.text.strip())
                        
                        for shape_idx, shape in enumerate(slide.shapes):
                            try:
                                if shape_idx == 0 and slide_content['title']:
                                    continue
                                
                                shape_info = {
                                    'type': shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type),
                                    'name': shape.name,
                                    'position': PPTStyleExtractor.extract_position(shape),
                                    'style': PPTStyleExtractor.extract_shape_style(shape)
                                }
                                
                                if hasattr(shape, 'text') and shape.text.strip():
                                    shape_info['text'] = TextCleaner.clean_text(shape.text.strip())
                                    
                                    if hasattr(shape, 'text_frame'):
                                        tf = shape.text_frame
                                        paragraphs = []
                                        for p_idx, para in enumerate(tf.paragraphs):
                                            para_info = {
                                                'text': TextCleaner.clean_text(para.text),
                                                'level': para.level,
                                                'paragraph_style': PPTStyleExtractor.extract_paragraph_style(para),
                                                'font_style': PPTStyleExtractor.extract_font_style(para.font) if para.font else {}
                                            }
                                            paragraphs.append(para_info)
                                        shape_info['paragraphs'] = paragraphs
                                    
                                    slide_content['shapes'].append(shape_info)
                                    
                                    if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE:
                                        shape_type = shape.shape_type
                                        shape_name = shape.name.lower() if shape.name else ''
                                        
                                        if 'picture' in shape_name or hasattr(shape, 'image'):
                                            try:
                                                image_info = {
                                                    'type': 'image',
                                                    'name': shape.name,
                                                    'width': None,
                                                    'height': None,
                                                    'description': '',
                                                    'data': None,
                                                    'is_placeholder': True
                                                }
                                                
                                                if hasattr(shape, 'width') and shape.width:
                                                    image_info['width'] = shape.width.pt
                                                if hasattr(shape, 'height') and shape.height:
                                                    image_info['height'] = shape.height.pt
                                                
                                                if hasattr(shape, 'alt_text') and shape.alt_text:
                                                    image_info['description'] = shape.alt_text
                                                
                                                if hasattr(shape, '_inline') and shape._inline:
                                                    try:
                                                        rId = shape._inline.docPr.get('id')
                                                        if rId:
                                                            rId_str = str(rId)
                                                            for img_name, img_data in image_map.items():
                                                                if img_data.get('rId') == rId_str or img_name:
                                                                    image_info['data'] = img_data.get('data')
                                                                    image_info['is_placeholder'] = False
                                                                    image_info['format'] = img_data.get('format')
                                                                    break
                                                    except:
                                                        pass
                                                
                                                slide_content['content_blocks'].append(image_info)
                                                slide_content['shapes'].append({
                                                    'type': 'PICTURE',
                                                    'name': shape.name,
                                                    'image_info': image_info
                                                })
                                            except Exception as pic_error:
                                                logger.warning(f'处理PPT图片失败: {str(pic_error)}')
                                                slide_content['content_blocks'].append({
                                                    'type': 'image',
                                                    'name': shape.name,
                                                    'description': '图片无法预览',
                                                    'is_placeholder': True
                                                })
                                        else:
                                            slide_content['content_blocks'].append({
                                                'type': 'shape',
                                                'text': TextCleaner.clean_text(shape.text)
                                            })
                                    elif hasattr(shape, 'table'):
                                        table_info = {
                                            'type': 'table',
                                            'rows': [],
                                            'first_row': [],
                                            'first_col': False,
                                            'last_row': False,
                                            'last_col': False,
                                            'band_row': False
                                        }
                                        table = shape.table
                                        for r_idx, row in enumerate(table.rows):
                                            row_data = []
                                            for cell in row.cells:
                                                row_data.append(cell.text.strip())
                                            table_info['rows'].append(row_data)
                                        slide_content['content_blocks'].append(table_info)
                                    else:
                                        slide_content['content_blocks'].append({
                                            'type': 'text',
                                            'text': shape.text.strip()
                                        })
                            except Exception as shape_error:
                                logger.warning(f'幻灯片 {i+1} 中形状 {shape_idx} 处理失败: {str(shape_error)}')
                                continue
                        
                        if not slide_content['title'] and slide_content['content_blocks']:
                            first_text = slide_content['content_blocks'][0].get('text', '') if slide_content['content_blocks'] else ''
                            if first_text:
                                slide_content['title'] = first_text[:50] + ('...' if len(first_text) > 50 else '')
                        
                        preview_options.append({
                            'type': 'slide',
                            'content': slide_content,
                            'name': f'幻灯片 {i + 1}',
                            'format': 'structured'
                        })
                        
                    except Exception as slide_error:
                        logger.warning(f'幻灯片 {i+1} 处理失败: {str(slide_error)}')
                        preview_options.append({
                            'type': 'text',
                            'content': f'幻灯片 {i+1} 处理失败: {str(slide_error)}',
                            'name': f'幻灯片 {i+1}'
                        })
                        continue
                
                if slide_count > slides_to_preview:
                    preview_options.append({
                        'type': 'text',
                        'content': f'文档包含{slide_count}张幻灯片，仅显示前{slides_to_preview}张',
                        'name': '提示信息'
                    })
                
            except ImportError:
                logger.warning('python-pptx未安装，尝试使用PyMuPDF')
                import fitz
                
                doc = fitz.open(file_path)
                for page_num in range(min(10, len(doc))):
                    page = doc.load_page(page_num)
                    text = page.get_text()
                    
                    if text.strip():
                        preview_options.append({
                            'type': 'text',
                            'content': text,
                            'name': f'幻灯片 {page_num + 1}',
                            'format': 'plain'
                        })
                doc.close()
                
                if len(doc) > 10:
                    preview_options.append({
                        'type': 'text',
                        'content': f'文档包含{len(doc)}张幻灯片，仅显示前10张',
                        'name': '提示信息'
                    })
            
            if not preview_options:
                preview_options.append({
                    'type': 'text',
                    'content': '演示文稿为空或无法提取内容',
                    'name': '提示信息'
                })
            
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'powerpoint',
                'file_format': 'pptx',
                'preview_options': preview_options,
                'images': image_map,
                'conversion_available': ['html', 'pdf', 'text']
            }
            
        except Exception as e:
            logger.error(f'新版PowerPoint文档处理失败: {str(e)}', exc_info=True)
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'powerpoint',
                'file_format': 'pptx',
                'preview_options': [{
                    'type': 'text',
                    'content': 'PowerPoint文档预览失败，建议下载后查看原文件。错误详情：' + str(e),
                    'name': '错误信息'
                }],
                'conversion_available': ['text'],
                'error': str(e)
            }
    
    @classmethod
    def _preview_old_powerpoint_file(cls, file_path: str, file_name: str) -> dict:
        """预览旧版PowerPoint文档(.ppt)"""
        try:
            # 方法1：使用PyMuPDF尝试提取文本（主要方法）
            try:
                import fitz
                
                doc = fitz.open(file_path)
                slide_contents = []
                doc_len = len(doc)
                
                # 提取每页文本
                for page_num in range(min(10, doc_len)):  # 最多预览10页
                    try:
                        page = doc.load_page(page_num)
                        text = page.get_text()
                        
                        if text.strip():
                            slide_contents.append({
                                'slide_number': page_num + 1,
                                'content': text
                            })
                    except Exception as page_error:
                        logger.warning(f'PyMuPDF处理幻灯片 {page_num+1} 失败: {str(page_error)}')
                        continue
                
                doc.close()
                
                # 构建预览结果
                preview_options = []
                
                # 幻灯片预览
                if slide_contents:
                    for slide in slide_contents:
                        preview_options.append({
                            'type': 'text',
                            'content': f"幻灯片 {slide['slide_number']}:\n{slide['content']}",
                            'name': f"幻灯片 {slide['slide_number']}",
                            'format': 'plain',
                            'note': '旧版PowerPoint文档，可能无法完全保留格式'
                        })
                else:
                    preview_options.append({
                        'type': 'text',
                        'content': '无法提取文档内容，请下载后查看',
                        'name': '提示信息'
                    })
                
                # 如果幻灯片数量超过10张，添加提示
                if doc_len > 10:
                    preview_options.append({
                        'type': 'text',
                        'content': f'文档包含{doc_len}张幻灯片，仅显示前10张',
                        'name': '提示信息'
                    })
                
                return {
                    'type': 'office_enhanced',
                    'name': file_name,
                    'office_type': 'powerpoint',
                    'file_format': 'ppt',
                    'preview_options': preview_options,
                    'conversion_available': ['text'],
                    'compatibility_note': '旧版PowerPoint文档，功能受限'
                }
                
            except Exception as fitz_error:
                logger.warning(f'PyMuPDF处理旧版PowerPoint失败: {str(fitz_error)}')
            
            # 方法2：尝试作为二进制文件读取基本信息
            try:
                file_size = os.path.getsize(file_path)
                file_info = f"文件名: {file_name}\n文件大小: {file_size / 1024:.2f} KB\n格式: 旧版PowerPoint (.ppt)"
                
                preview_options = [{
                    'type': 'text',
                    'content': f"{file_info}\n\n系统无法直接解析此旧版PowerPoint文件的内容。请下载后使用Microsoft PowerPoint查看完整内容。",
                    'name': '文件信息',
                    'note': '使用基本文件信息作为备选方案'
                }]
                
                return {
                    'type': 'office_enhanced',
                    'name': file_name,
                    'office_type': 'powerpoint',
                    'file_format': 'ppt',
                    'preview_options': preview_options,
                    'conversion_available': ['text'],
                    'compatibility_note': '旧版PowerPoint文档，仅能显示基本信息',
                    'error': '无法提取内容，但提供基本文件信息'
                }
            except Exception as binary_error:
                logger.error(f'基本信息提取失败: {str(binary_error)}')
            
            # 如果所有方法都失败，返回友好的错误信息
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'powerpoint',
                'file_format': 'ppt',
                'preview_options': [{
                    'type': 'text',
                    'content': '旧版PowerPoint文档预览失败，可能是文件格式不兼容或文件损坏。建议下载后使用Microsoft PowerPoint查看。',
                    'name': '提示信息',
                    'note': '系统尝试了多种提取方法但都失败了'
                }],
                'conversion_available': ['text'],
                'error': '无法提取文档内容'
            }
            
        except Exception as e:
            logger.error(f'旧版PowerPoint文档处理失败: {str(e)}', exc_info=True)
            # 即使处理失败，也返回基本信息，而不是抛出异常
            return {
                'type': 'office_enhanced',
                'name': file_name,
                'office_type': 'powerpoint',
                'file_format': 'ppt',
                'preview_options': [{
                    'type': 'text',
                    'content': 'PowerPoint文档预览失败，建议下载后查看原文件。错误详情：' + str(e),
                    'name': '错误信息'
                }],
                'conversion_available': ['text'],
                'error': str(e)
            }
    
    @classmethod
    def _convert_to_image(cls, file_path: str, file_name: str, office_type: str, file_ext: str) -> dict:
        """将Office文档转换为图片格式"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                if office_type == 'word':
                    return cls._convert_word_to_image(file_path, file_name, file_ext, temp_dir)
                elif office_type == 'excel':
                    return cls._convert_excel_to_image(file_path, file_name, file_ext, temp_dir)
                elif office_type == 'powerpoint':
                    return cls._convert_powerpoint_to_image(file_path, file_name, file_ext, temp_dir)
                else:
                    raise ValueError(f'不支持的Office类型: {office_type}')
                    
        except Exception as e:
            logger.error(f'图片转换失败: {str(e)}', exc_info=True)
            return cls._create_error_response(file_name, f'图片转换失败: {str(e)}')
    
    @classmethod
    def _convert_word_to_image(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将Word文档转换为图片"""
        try:
            # 使用PyMuPDF将Word文档转换为图片
            import fitz
            
            doc = fitz.open(file_path)
            
            # 创建图片列表
            image_list = []
            
            # 将每页转换为图片
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # 设置图片缩放比例
                mat = fitz.Matrix(2.0, 2.0)  # 2倍缩放，提高清晰度
                pix = page.get_pixmap(matrix=mat)
                
                # 保存图片
                image_path = os.path.join(temp_dir, f"page_{page_num + 1}.png")
                pix.save(image_path)
                
                # 读取图片内容
                with open(image_path, 'rb') as f:
                    image_data = f.read()
                
                image_list.append({
                    'page': page_num + 1,
                    'data': image_data,
                    'format': 'png',
                    'size': len(image_data)
                })
            
            doc.close()
            
            return {
                'type': 'image',
                'name': f"{os.path.splitext(file_name)[0]}.png",
                'images': image_list,
                'total_pages': len(image_list),
                'file_format': 'png'
            }
                
        except Exception as e:
            raise Exception(f'Word文档图片转换失败: {str(e)}')
    
    @classmethod
    def _convert_excel_to_image(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将Excel文档转换为图片"""
        try:
            # 对于Excel文档，先转换为HTML，再转换为图片
            # 这里使用简单的文本表示，实际项目中可以使用更复杂的转换方法
            
            # 读取Excel内容
            text_content = []
            
            if file_ext == '.xlsx':
                from openpyxl import load_workbook
                
                wb = load_workbook(file_path, data_only=True)
                
                for sheet_name in wb.sheetnames[:3]:  # 最多处理3个工作表
                    ws = wb[sheet_name]
                    
                    text_content.append(f'=== {sheet_name} ===')
                    
                    # 获取数据范围
                    max_row = min(10, ws.max_row)  # 最多10行
                    max_col = min(5, ws.max_column)  # 最多5列
                    
                    for row in range(1, max_row + 1):
                        row_data = []
                        for col in range(1, max_col + 1):
                            cell_value = ws.cell(row=row, column=col).value
                            if cell_value is None:
                                cell_value = ''
                            row_data.append(str(cell_value))
                        text_content.append(' | '.join(row_data))
                    
                    text_content.append('')  # 空行分隔
                
                wb.close()
                
            else:  # .xls
                import xlrd
                
                wb = xlrd.open_workbook(file_path)
                
                for sheet_index in range(min(3, wb.nsheets)):  # 最多处理3个工作表
                    ws = wb.sheet_by_index(sheet_index)
                    
                    text_content.append(f'=== {ws.name} ===')
                    
                    # 获取数据范围
                    max_row = min(10, ws.nrows)  # 最多10行
                    max_col = min(5, ws.ncols)  # 最多5列
                    
                    for row in range(max_row):
                        row_data = []
                        for col in range(max_col):
                            cell_value = ws.cell_value(row, col)
                            if cell_value is None:
                                cell_value = ''
                            row_data.append(str(cell_value))
                        text_content.append(' | '.join(row_data))
                    
                    text_content.append('')  # 空行分隔
                
                wb.release_resources()
            
            if not text_content:
                text_content = ['工作簿为空']
            
            # 创建简单的文本图片表示
            # 在实际项目中，可以使用PIL等库生成真正的图片
            return {
                'type': 'image',
                'name': f"{os.path.splitext(file_name)[0]}.png",
                'content': '\n'.join(text_content),
                'file_format': 'text_image',
                'note': 'Excel文档图片转换，实际项目中可使用更复杂的图片生成方法'
            }
                
        except Exception as e:
            raise Exception(f'Excel文档图片转换失败: {str(e)}')
    
    @classmethod
    def _convert_powerpoint_to_image(cls, file_path: str, file_name: str, file_ext: str, temp_dir: str) -> dict:
        """将PowerPoint文档转换为图片"""
        try:
            # 使用PyMuPDF将PowerPoint文档转换为图片
            import fitz
            
            doc = fitz.open(file_path)
            
            # 创建图片列表
            image_list = []
            
            # 将每页转换为图片
            for page_num in range(min(10, len(doc))):  # 最多处理10张幻灯片
                page = doc.load_page(page_num)
                
                # 设置图片缩放比例
                mat = fitz.Matrix(2.0, 2.0)  # 2倍缩放，提高清晰度
                pix = page.get_pixmap(matrix=mat)
                
                # 保存图片
                image_path = os.path.join(temp_dir, f"slide_{page_num + 1}.png")
                pix.save(image_path)
                
                # 读取图片内容
                with open(image_path, 'rb') as f:
                    image_data = f.read()
                
                image_list.append({
                    'slide': page_num + 1,
                    'data': image_data,
                    'format': 'png',
                    'size': len(image_data)
                })
            
            doc.close()
            
            return {
                'type': 'image',
                'name': f"{os.path.splitext(file_name)[0]}.png",
                'images': image_list,
                'total_slides': len(image_list),
                'file_format': 'png'
            }
                
        except Exception as e:
            raise Exception(f'PowerPoint文档图片转换失败: {str(e)}')